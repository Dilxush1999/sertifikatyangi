import os
import sys
import datetime
import json
import sqlite3
import datetime
import threading
import asyncio
from concurrent.futures import ThreadPoolExecutor
from telegram import Update, ReplyKeyboardMarkup, ReplyKeyboardRemove, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    ApplicationBuilder, CommandHandler, MessageHandler, ConversationHandler, ContextTypes, filters, CallbackQueryHandler
)
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from PIL import Image
import img2pdf
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import requests
from pptx.enum.text import PP_ALIGN
from tenacity import retry, stop_after_attempt, wait_exponential
import csv
import io
import img2pdf
# Bosh importlarga qo'shing:
from pdf2image import convert_from_path
from PIL import ImageDraw, ImageFont
import logging

from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from PIL import Image
import shutil
from contextlib import contextmanager
import time
import qrcode

# .exe fayli uchun dinamik yoâ€˜l
if getattr(sys, 'frozen', False):
    BASE_DIR = os.path.dirname(sys.executable)
else:
    BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# SQLite uchun datetime adapter va converter
def adapt_datetime(dt):
    return dt.isoformat()

def convert_datetime(s):
    return datetime.datetime.fromisoformat(s.decode('utf-8'))

sqlite3.register_adapter(datetime.datetime, adapt_datetime)
sqlite3.register_converter("datetime", convert_datetime)

db_lock = threading.Lock()
file_lock = threading.Lock()


TEMP_PATH = os.path.join(BASE_DIR, 'temp')
FONT_DIR = os.path.join(BASE_DIR, 'Font')
SERTIFIKAT_PATH = os.path.join(BASE_DIR, 'Sertifikat')
TAKLIFNOMA_PATH = os.path.join(BASE_DIR, 'Taklifnoma')
SHABLON_PATH = os.path.join(BASE_DIR, 'Shablonlar')
DIPLOM_PATH = os.path.join(BASE_DIR, 'Diplom')
IMAGE_PATH = os.path.join(BASE_DIR, 'Image')
USERS_PATH = os.path.join(BASE_DIR, 'Users')
SAMPLE_IMAGE = os.path.join(IMAGE_PATH, 'namuna.jpg')
INFOTEXT_IMAGE = os.path.join(IMAGE_PATH, 'infotext.jpg')
TAKLIFNOMA_IMAGE = os.path.join(IMAGE_PATH, 'taklifnoma_info.jpg')
DIPLOM_MATNI_IMAGE = os.path.join(IMAGE_PATH, 'diplom_matni.jpg')
DIPLOM_NAMUNA_IMAGE = os.path.join(IMAGE_PATH, 'diplom_namuna.jpg')
DIPLOM_NAMUNA_IMAGE1 = os.path.join(IMAGE_PATH, 'diplom_namuna1.jpg')
DIPLOM_NAMUNA_IMAGE2 = os.path.join(IMAGE_PATH, 'diplom_namuna2.jpg')
DIPLOM_NAMUNA_IMAGE3 = os.path.join(IMAGE_PATH, 'diplom_namuna3.jpg')
DIPLOM_NAMUNA_IMAGE4 = os.path.join(IMAGE_PATH, 'diplom_namuna4.jpg')
SAMPLE_IMAGE1 = os.path.join(IMAGE_PATH, 'namuna1.jpg')
SAMPLE_IMAGE2 = os.path.join(IMAGE_PATH, 'namuna2.jpg')
SAMPLE_IMAGE3 = os.path.join(IMAGE_PATH, 'namuna3.jpg')
SAMPLE_IMAGE4 = os.path.join(IMAGE_PATH, 'namuna4.jpg')
TAKLIFNOMA_SAMPLE_IMAGE = os.path.join(IMAGE_PATH, 'taklifnoma_namuna.jpg')
TAKLIFNOMA_SAMPLE_IMAGE1 = os.path.join(IMAGE_PATH, 'taklifnoma_namuna1.jpg')
TAKLIFNOMA_SAMPLE_IMAGE2 = os.path.join(IMAGE_PATH, 'taklifnoma_namuna2.jpg')
SHABLON_IMAGE = os.path.join(IMAGE_PATH, 'shablon.jpg')
SHABLON_IMAGE1 = os.path.join(IMAGE_PATH, 'shablon1.jpg')
SHABLON_IMAGE2 = os.path.join(IMAGE_PATH, 'shablon2.jpg')
SHABLON_NOMI_IMAGE = os.path.join(IMAGE_PATH, 'shablon_nomi.jpg')
SHABLON_MATNI_IMAGE = os.path.join(IMAGE_PATH, 'shablon_matni.jpg')
SHRIFT_IMAGE = os.path.join(IMAGE_PATH, 'shrift.jpg')
CONFIG_PATH = os.path.join(BASE_DIR, 'config.json')
DB_PATH = os.path.join(BASE_DIR, 'bot_db.sqlite')

for path in [TEMP_PATH, USERS_PATH, SERTIFIKAT_PATH, TAKLIFNOMA_PATH, SHABLON_PATH, DIPLOM_PATH, IMAGE_PATH]:
    if not os.path.exists(path):
        os.makedirs(path)

START, DOCUMENT_TYPE, TEMPLATE, TAQDIRLANUVCHI, SHRIFT1, TAQDIRLOVCHI, SERTIFIKAT_MATNI, SHRIFT2, SANA, BALANCE, PAYMENT_METHOD, CARD_PAYMENT, UPLOAD_CHECK, COMMENT, ADMIN_PANEL, ADMIN_ACTION, ADMIN_USER, ADMIN_TOPUP, INFO_TEXT, MANZIL_VA_EGA, SHABLON, SHABLON_NOMI, SHABLON_SHRIFT1, SHABLON_MATNI, SHABLON_SHRIFT2, SHABLON_SHRIFT3, TAKLIFNOMA_SHRIFT1, TAKLIFNOMA_SHRIFT2, PDF_CONFIRM, DIPLOM_MATNI, ADMIN_MESSAGE_TYPE, ADMIN_MESSAGE_CONTENT, ADMIN_MESSAGE_RECIPIENT, CONTACT, CONTACT_MESSAGE, CONFIG_PRICE, SET_NEW_PRICE, ADMIN_FOYDALANUVCHI, QR_CODE = range(39)
SERTIFIKAT_TAGS = ["{taqdirlangan}", "{taqdirlovchi}", "{sertifikat_matni}", "{sana}", "{qr_code}"]
TAKLIFNOMA_TAGS = ["{info_text}", "{manzil_va_ega}"]
DIPLOM_TAGS = ["{taqdirlangan}", "{taqdirlovchi}", "{diplom_matni}", "{sana}", "{qr_code}"]
SHABLON_TAGS = ["{shablon_nomi}", "{shablon_matni}", "{sana}", "{taqdirlovchi}", "{qr_code}"]

FONTS = [
    "Times New Roman", "Bodoni MT", "Algerian", "Castellar", "Comic Sans MS",
    "Constantia", "Elephant", "Freestyle Script", "Harlow Solid Italic", "Impact",
    "Informal Roman", "Lucida Calligraphy", "Lucida Handwriting", "Yu Gothic", "Verdana",
    "Monotype Corsiva", "Palatino Linotype", "Pristina", "Old English Text MT", "Niagara Engraved",
    "Vivaldi", "Matura MT Script Capitals", "Ink Free", "Kristen ITC", "Gabriola",
    "Eras Demi ITC", "Bell MT", "Tempus Sans ITC", "Harrington", "Jokerman",
    "MS Mincho", "Eras Medium ITC", "Papyrus", "Stencil", "Script MT Bold",
    "MV Boli", "Calibri Light", "Ravie", "Curlz MT", "Magneto"
]


logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO,
    handlers=[
        logging.FileHandler('bot.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Global navbat tizimi
document_queue = asyncio.Queue()
queue_worker_running = False

def load_config():
    try:
        with open('config.json', 'r') as f:
            config = json.load(f)
        required_keys = ['BOT_TOKEN', 'CERTIFICATE_COST', 'TAKLIFNOMA_COST', 'SHABLON_COST', 'DIPLOM_COST', 'PDF_COST', 'REFERRAL_BONUS']
        missing_keys = [key for key in required_keys if key not in config]
        if missing_keys:
            raise KeyError(f"Config faylida quyidagi kalitlar yoâ€˜q: {missing_keys}")
        return config
    except Exception as e:
        logger.error(f"Konfiguratsiya faylini yuklashda xato: {str(e)}")
        raise

@contextmanager
def get_db_connection():
    """Ma'lumotlar bazasiga ulanishni boshqarish"""
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=1, max=10))
    def connect():
        conn = sqlite3.connect(DB_PATH, detect_types=sqlite3.PARSE_DECLTYPES)
        conn.row_factory = sqlite3.Row
        return conn

    try:
        conn = connect()
        yield conn
    except sqlite3.Error as e:
        logger.error(f"DB ulanish xatosi: {str(e)}")
        raise
    finally:
        if conn:
            conn.close()
            logger.debug("DB ulanishi yopildi")

def init_db():
    try:
        with get_db_connection() as conn:
            c = conn.cursor()
            c.execute('''CREATE TABLE IF NOT EXISTS users (
                user_id TEXT PRIMARY KEY,
                username TEXT NOT NULL,
                balance INTEGER NOT NULL CHECK (balance >= 0),
                taqdirlangan TEXT,
                taqdirlovchi TEXT,
                sertifikat_matni TEXT,
                sana TEXT,
                info_text TEXT,
                manzil_va_ega TEXT,
                shablon_nomi TEXT,
                shablon_matni TEXT,
                referrals INTEGER NOT NULL DEFAULT 0 CHECK (referrals >= 0),
                phone_number TEXT,
                referrals_ids TEXT DEFAULT '[]',
                diplom_matni TEXT,
                result_count INTEGER NOT NULL DEFAULT 0 CHECK (result_count >= 0),
                is_blocked INTEGER NOT NULL DEFAULT 0  -- Yangi ustun
            )''')
            c.execute('''CREATE TABLE IF NOT EXISTS transactions (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id TEXT NOT NULL,
                amount INTEGER NOT NULL,
                transaction_type TEXT NOT NULL,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )''')
            c.execute('''CREATE TABLE IF NOT EXISTS cache (
                key TEXT PRIMARY KEY,
                value TEXT NOT NULL,
                created_at TEXT NOT NULL,
                expires_at TEXT
            )''')
            conn.commit()
            logger.info("Ma'lumotlar bazasi muvaffaqiyatli ishga tushirildi")
    except sqlite3.Error as e:
        logger.error(f"Ma'lumotlar bazasini ishga tushirishda xato: {str(e)}")
        raise

def log_transaction(user_id, amount, transaction_type):
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("INSERT INTO transactions (user_id, amount, transaction_type) VALUES (?, ?, ?)",
                  (str(user_id), amount, transaction_type))
        conn.commit()

def save_to_cache(key, value, ttl=3600):
    expires_at = datetime.datetime.now() + datetime.timedelta(seconds=ttl)
    created_at = datetime.datetime.now()
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("INSERT OR REPLACE INTO cache (key, value, created_at, expires_at) VALUES (?, ?, ?, ?)",
                  (key, value, created_at.isoformat(), expires_at.isoformat()))
        conn.commit()

def get_from_cache(key):
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("DELETE FROM cache WHERE expires_at IS NOT NULL AND expires_at < ?",
                  (datetime.datetime.now().isoformat(),))
        c.execute("SELECT value FROM cache WHERE key = ?", (key,))
        result = c.fetchone()
        return result['value'] if result else None

def delete_from_cache(key):
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("DELETE FROM cache WHERE key = ?", (key,))
        conn.commit()

def load_user_data(user_id):
    logger.debug(f"Loading user data for user_id: {user_id}")
    cache_key = f"user:{user_id}"
    cached_data = get_from_cache(cache_key)
    if cached_data:
        logger.debug(f"Cache hit for {user_id}")
        return json.loads(cached_data)

    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("SELECT * FROM users WHERE user_id = ?", (str(user_id),))
        user = c.fetchone()
        if user:
            user_data = {
                "username": user['username'],
                "balance": user['balance'],
                "taqdirlangan": user['taqdirlangan'],
                "taqdirlovchi": user['taqdirlovchi'],
                "sertifikat_matni": user['sertifikat_matni'],
                "sana": user['sana'],
                "info_text": user['info_text'],
                "manzil_va_ega": user['manzil_va_ega'],
                "shablon_nomi": user['shablon_nomi'],
                "shablon_matni": user['shablon_matni'],
                "referrals": user['referrals'],
                "referrals_ids": json.loads(user['referrals_ids']),
                "diplom_matni": user['diplom_matni'],
                "result_count": user['result_count'],
                "is_blocked": user['is_blocked']  # Yangi maydon
            }
            save_to_cache(cache_key, json.dumps(user_data))
            return user_data
        default_data = {
            "balance": 5000,
            "referrals": 0,
            "referrals_ids": [],
            "taqdirlangan": "Alisherov Farid",
            "taqdirlovchi": "Bahodirov Dilxush",
            "sertifikat_matni": (
                "Siz bizning Telegram Bot 3-oylik kursimizni 1 - o'rinda muvaffaqiyatli tamomlaganingiz "
                "uchun ushbu sertifikat bilan taqdirlanasiz va sizga jamoamiz nomidan lutfan tashakkur bildiramiz!"
            ),
            "info_text": (
                "Hurmatli AZIZ MEHMONIMIZ\n"
                "Siz va oila a'zolaringizni 2025-yil 16-fevral kuni soat 18:00 da aziz farzandimiz\n"
                "FARIDJON va SHUKRONAXON larning\n"
                "Nikoh to'yi munosabai bilan yoziladigan dasturxonimizga lutgan taklif etamiz!"
            ),
            "manzil_va_ega": (
                "Hurmat bilan Bahodirovlar oilasi\n"
                "Manzil: Dehqonobod tumani Mamat ota to'yxonasi"
            ),
            "shablon_nomi": "TASHAKKURNOMA",
            "shablon_matni": (
                "Qashqadaryo viloyati Dehqonobod tumani 87â€“sonli umumta'lim maktabining 5-'A' sinf oâ€˜quvchisi "
                "Alisherov Farid 2024 â€“ 2025 O'quv yilida a'lo va yaxshi baholari hamda namunali xulqi uchun "
                "MAQTOV YORLIG'I bilan taqdirlanadi"
            ),
            "diplom_matni": "Siz bizning kursimizni muvaffaqiyatli tamomlaganingiz uchun ushbu diplom bilan taqdirlanasiz!",
            "result_count": 0,
            "is_blocked": 0  # Yangi maydon
        }
        save_to_cache(cache_key, json.dumps(default_data))
        return default_data

def save_user_data(user_id, data):
    logger.debug(f"Saving user data for user_id: {user_id}, data: {data}")
    if data.get('balance', 5000) < 0:
        raise ValueError("Balans manfiy boâ€˜lishi mumkin emas")
    if data.get('result_count', 0) < 0:
        raise ValueError("Natijalar soni manfiy boâ€˜lishi mumkin emas")
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute('''INSERT INTO users (
            user_id, username, balance, taqdirlangan, taqdirlovchi, sertifikat_matni, sana, 
            info_text, manzil_va_ega, shablon_nomi, shablon_matni, referrals, referrals_ids, 
            diplom_matni, result_count, is_blocked
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ON CONFLICT (user_id) DO UPDATE SET
            username = EXCLUDED.username,
            balance = EXCLUDED.balance,
            taqdirlangan = EXCLUDED.taqdirlangan,
            taqdirlovchi = EXCLUDED.taqdirlovchi,
            sertifikat_matni = EXCLUDED.sertifikat_matni,
            sana = EXCLUDED.sana,
            info_text = EXCLUDED.info_text,
            manzil_va_ega = EXCLUDED.manzil_va_ega,
            shablon_nomi = EXCLUDED.shablon_nomi,
            shablon_matni = EXCLUDED.shablon_matni,
            referrals = EXCLUDED.referrals,
            referrals_ids = EXCLUDED.referrals_ids,
            diplom_matni = EXCLUDED.diplom_matni,
            result_count = EXCLUDED.result_count,
            is_blocked = EXCLUDED.is_blocked''',
            (
                str(user_id), 
                data.get('username', 'Nomaâ€˜lum'), 
                data.get('balance', 5000),
                data.get('taqdirlangan'), 
                data.get('taqdirlovchi'), 
                data.get('sertifikat_matni'), 
                data.get('sana'),
                data.get('info_text'), 
                data.get('ãƒžãƒ³zil_va_ega'), 
                data.get('shablon_nomi'), 
                data.get('shablon_matni'), 
                data.get('referrals', 0),
                json.dumps(data.get('referrals_ids', [])),
                data.get('diplom_matni'),
                data.get('result_count', 0),
                data.get('is_blocked', 0)
            )
        )
        conn.commit()
    cache_key = f"user:{user_id}"
    save_to_cache(cache_key, json.dumps(data))

def get_all_text_from_presentation(prs):
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
            if shape.shape_type == 6:
                for subshape in shape.shapes:
                    if hasattr(subshape, "text"):
                        texts.append(subshape.text)
    return texts

@contextmanager
def temp_file_manager(user_id, timestamp):
    """Vaqtinchalik fayllarni boshqarish uchun kontekst menejeri"""
    temp_dir = os.path.join(TEMP_PATH, f"user_{user_id}_{timestamp}")
    os.makedirs(temp_dir, exist_ok=True)
    temp_files = []
    try:
        yield temp_dir, temp_dir
    finally:
        for file_path in temp_files:
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    logger.debug(f"Fayl oâ€˜chirildi: {file_path}")
                except OSError as e:
                    logger.error(f"Faylni oâ€˜chirishda xato: {file_path}, {str(e)}")
        if os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
                logger.info(f"Temp papka oâ€˜chirildi: {temp_dir}")
            except Exception as e:
                logger.error(f"Temp papkani oâ€˜chirishda xato: {temp_dir}, {str(e)}")

def check_template(template_path, tags):
    if not os.path.exists(template_path):
        return False, f"âŒ {template_path} fayli topilmadi.", False, []
    try:
        prs = Presentation(template_path)
        texts = get_all_text_from_presentation(prs)
        missing = []
        present_tags = []
        for tag in tags:
            if any(tag in text for text in texts):
                present_tags.append(tag)
            else:
                missing.append(tag)
        has_date_tag = "{sana}" in present_tags if "{sana}" in tags else False
        has_qr_code_tag = "{qr_code}" in present_tags if "{qr_code}" in tags else False
        if not present_tags:
            return False, f"âŒ {template_path} faylida hech qanday teg topilmadi.", has_date_tag, present_tags, has_qr_code_tag
        return True, f"âœ… {template_path} tayyor!", has_date_tag, present_tags, has_qr_code_tag
    except Exception as e:
        return False, f"âŒ {template_path} faylini oâ€˜qishda xatolik: {str(e)}", False, [], False

def save_context_data(user_id, context_data):
    with file_lock:
        with open(os.path.join(USERS_PATH, f"context_{user_id}.json"), "w", encoding="utf-8") as f:
            json.dump(context_data, f, ensure_ascii=False, indent=4)

def load_context_data(user_id):
    with file_lock:
        try:
            with open(os.path.join(USERS_PATH, f"context_{user_id}.json"), "r", encoding="utf-8") as f:
                return json.load(f)
        except FileNotFoundError:
            return {}

def get_templates(path):
    files = [f for f in os.listdir(path) if f.endswith('.pptx') and f.split('.')[0].isdigit()]
    files.sort(key=lambda x: int(x.split('.')[0]))
    return files

def create_template_keyboard(templates=None):
    if templates:
        keyboard = []
        for i in range(0, len(templates), 5):
            row = [str(int(f.split('.')[0])) for f in templates[i:i+5]]
            keyboard.append(row)
        return keyboard
    return [
        ["1", "2", "3", "4", "5"],
        ["6", "7", "8", "9", "10"],
        ["11", "12", "13", "14", "15"]
    ]

def create_font_keyboard():
    keyboard = []
    for i in range(0, len(FONTS), 5):
        row = [str(i + 1) for i in range(i, min(i + 5, len(FONTS)))]
        keyboard.append(row)
    return keyboard

def validate_date(date_text):
    try:
        datetime.datetime.strptime(date_text, "%d.%m.%Y")
        return True
    except ValueError:
        return False

def check_fonts():
    """Font papkasidagi shrift fayllarini tekshirish"""
    global FONT_DIR
    FONT_DIR = os.path.join(BASE_DIR, 'Font')
    if not os.path.exists(FONT_DIR):
        logger.warning(f"Font papkasi topilmadi: {FONT_DIR}. Default shrift ishlatiladi.")
        return

    missing_fonts = []
    for i in range(1, 41):
        font_path = os.path.join(FONT_DIR, f"{i}.ttf")
        if not os.path.exists(font_path):
            missing_fonts.append(f"{i}.ttf")
    
    if missing_fonts:
        logger.warning(f"Ba'zi shrift fayllari topilmadi: {missing_fonts}. Default shrift ishlatiladi.")
    else:
        logger.info("Barcha shrift fayllari tekshirildi va topildi")

async def convert_pptx_to_jpg(pptx_path, output_path, context):
    """
    PPTX faylini JPG formatiga aylantiradi, placeholder matnlarni foydalanuvchi ma'lumotlari bilan
    almashtirib, barcha matnlarni saqlaydi va QR kodni qo'shadi.
    """
    try:
        logger.info(f"PPTX to JPG conversion started for {pptx_path}")
        prs = Presentation(pptx_path)
        if not prs.slides:
            logger.error("PPTX faylida slaydlar topilmadi")
            return False

        slide = prs.slides[0]
        # Yuqori sifat uchun o'lchamni 2x kattalashtirish
        slide_width = int(prs.slide_width.pt * 2)
        slide_height = int(prs.slide_height.pt * 2)
        img = Image.new('RGB', (slide_width, slide_height), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Shrift konfiguratsiyasi
        font_config = {
            'taqdirlanuvchi_font': {'font': context.user_data.get('taqdirlanuvchi_font', 'Arial'), 'size': 24},
            'other_font': {'font': context.user_data.get('other_font', 'Arial'), 'size': 16},
            'sana_font': {'font': context.user_data.get('other_font', 'Arial'), 'size': 15},
            'info_text_font': {'font': context.user_data.get('info_text_font', 'Times New Roman'), 'size': 26},
            'manzil_va_ega_font': {'font': context.user_data.get('manzil_va_ega_font', 'Times New Roman'), 'size': 22},
            'shablon_nomi_font': {'font': context.user_data.get('shablon_nomi_font', 'Arial'), 'size': 26},
            'shablon_matni_font': {'font': context.user_data.get('shablon_matni_font', 'Arial'), 'size': 18},
            'diplom_matni_font': {'font': context.user_data.get('other_font', 'Arial'), 'size': 16},
            'default_font': {'font': 'Times New Roman', 'size': 22}  # Standart shrift oddiy matnlar uchun
        }

        # QR kodni tayyorlash
        qr_code_path = None
        if context.user_data.get('has_qr_code_tag') and context.user_data.get('qr_code_data'):
            qr_code_path = os.path.join(TEMP_PATH, f"qr_code_{context.user_data.get('user_id')}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.png")
            create_qr_code(context.user_data['qr_code_data'], qr_code_path)
            logger.info(f"QR kod yaratildi: {qr_code_path}")

        def wrap_text(text, font, max_width):
            """Matnni max_width ga mos ravishda qatorlarga bo'ladi."""
            lines = []
            words = text.split()
            current_line = []
            current_width = 0

            for word in words:
                test_line = ' '.join(current_line + [word])
                text_bbox = draw.textbbox((0, 0), test_line, font=font)
                text_width = text_bbox[2] - text_bbox[0]

                if text_width <= max_width:
                    current_line.append(word)
                    current_width = text_width
                else:
                    if current_line:
                        lines.append(' '.join(current_line))
                    current_line = [word]
                    text_bbox = draw.textbbox((0, 0), word, font=font)
                    current_width = text_bbox[2] - text_bbox[0]

            if current_line:
                lines.append(' '.join(current_line))

            return lines

        for shape in slide.shapes:
            logger.debug(f"Processing shape: {shape.shape_type}")
            if shape.has_text_frame:
                text = shape.text.strip()
                if not text:
                    continue
                logger.debug(f"Processing text: {text}")

                font_key = None
                font_size = None
                font_name = 'Arial'

                # Placeholder matnlarni almashtirish
                if "{taqdirlangan}" in text:
                    font_key = 'taqdirlanuvchi_font'
                    text = context.user_data.get('taqdirlangan', text)
                elif "{sertifikat_matni}" in text:
                    font_key = 'other_font'
                    text = context.user_data.get('sertifikat_matni', text)
                elif "{diplom_matni}" in text:
                    font_key = 'diplom_matni_font'
                    text = context.user_data.get('diplom_matni', text)
                elif "{taqdirlovchi}" in text:
                    font_key = 'other_font'
                    text = context.user_data.get('taqdirlovchi', text)
                elif "{sana}" in text:
                    font_key = 'sana_font'
                    text = context.user_data.get('sana', text)
                elif "{info_text}" in text:
                    font_key = 'info_text_font'
                    text = context.user_data.get('info_text', text)
                elif "{manzil_va_ega}" in text:
                    font_key = 'manzil_va_ega_font'
                    text = context.user_data.get('manzil_va_ega', text)
                elif "{shablon_nomi}" in text:
                    font_key = 'shablon_nomi_font'
                    text = context.user_data.get('shablon_nomi', text)
                elif "{shablon_matni}" in text:
                    font_key = 'shablon_matni_font'
                    text = context.user_data.get('shablon_matni', text)
                elif "{qr_code}" in text:
                    if qr_code_path and os.path.exists(qr_code_path):
                        qr_img = Image.open(qr_code_path)
                        qr_img = qr_img.resize((200, 200), Image.Resampling.LANCZOS)  # QR kod o'lchamini 2x masshtab
                        img.paste(qr_img, (int(shape.left.pt * 2), int(shape.top.pt * 2)))
                        logger.info(f"QR kod joylashtirildi: ({shape.left.pt * 2}, {shape.top.pt * 2})")
                        continue  # QR kod joylashtirilgandan so'ng matnni chizishni o'tkazib yuboramiz
                    else:
                        logger.warning("QR kod fayli topilmadi yoki yaratilmadi")
                        continue
                else:
                    font_key = 'default_font'  # Teg bo'lmagan matnlar uchun standart shrift

                # Shrift va o'lchamni font_config dan olish (2x masshtab)
                font_name = font_config[font_key]['font']
                font_size = font_config[font_key]['size'] * 2
                logger.info(f"Shrift tanlandi: {font_key} = {font_name}, o'lchami: {font_size} pt")

                # Shrift faylini yuklash
                font_index = None
                try:
                    font_index = FONTS.index(font_name)
                    logger.info(f"Shrift indeksi: {font_name}, indeks: {font_index}")
                except ValueError:
                    logger.warning(f"Shrift topilmadi: {font_name}. Standart shrift ishlatiladi.")
                    font_name = 'Arial'
                    font_index = None

                font_path = os.path.join(FONT_DIR, f"{font_index + 1}.ttf") if font_index is not None else None
                try:
                    if font_path and os.path.exists(font_path):
                        font = ImageFont.truetype(font_path, font_size)
                        logger.info(f"Shrift yuklandi: {font_path}, o'lchami: {font_size}")
                    else:
                        # Use default font with proper size
                        try:
                            # DejaVu Sans yoki boshqa TTF fontni ishlatishga harakat qilish
                            default_fonts = [
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                                "/usr/share/fonts/TTF/arial.ttf",
                                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
                            ]
                            font_loaded = False
                            for default_font_path in default_fonts:
                                if os.path.exists(default_font_path):
                                    font = ImageFont.truetype(default_font_path, font_size)
                                    logger.info(f"Default TTF shrift ishlatildi: {default_font_path}, o'lchami: {font_size}")
                                    font_loaded = True
                                    break
                            
                            if not font_loaded:
                                font = ImageFont.load_default()
                                logger.warning(f"Oddiy default shrift ishlatildi, o'lchami: {font_size}")
                        except Exception as e:
                            font = ImageFont.load_default()
                            logger.warning(f"Default shrift yuklashda xato: {str(e)}")
                except Exception as e:
                    logger.error(f"Shriftni yuklashda xato: {font_path}, xato: {str(e)}")
                    try:
                        # Fallback shrift
                        default_fonts = [
                            "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                            "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
                        ]
                        font_loaded = False
                        for default_font_path in default_fonts:
                            if os.path.exists(default_font_path):
                                font = ImageFont.truetype(default_font_path, font_size)
                                logger.info(f"Fallback TTF shrift ishlatildi: {default_font_path}")
                                font_loaded = True
                                break
                        
                        if not font_loaded:
                            font = ImageFont.load_default()
                    except:
                        font = ImageFont.load_default()

                # Text box o'lchamlarini olish (2x masshtab)
                shape_width = (shape.width.pt * 2) if shape.width else slide_width
                shape_height = (shape.height.pt * 2) if shape.height else slide_height

                # Matnni qatorlarga bo'lish
                lines = wrap_text(text, font, shape_width)
                logger.debug(f"Matn {len(lines)} qatorga bo'lindi: {lines}")

                # Qatorlarning umumiy balandligini hisoblash
                line_height = draw.textbbox((0, 0), "Ay", font=font)[3] - draw.textbbox((0, 0), "Ay", font=font)[1]
                total_height = len(lines) * line_height

                # Agar matn balandligi text box dan oshsa, shrift o'lchamini kichraytirish
                while total_height > shape_height and font_size > 8:
                    font_size -= 1
                    try:
                        if font_path and os.path.exists(font_path):
                            font = ImageFont.truetype(font_path, font_size)
                        else:
                            # Default shriftni qayta yuklash
                            default_fonts = [
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"
                            ]
                            font_loaded = False
                            for default_font_path in default_fonts:
                                if os.path.exists(default_font_path):
                                    font = ImageFont.truetype(default_font_path, font_size)
                                    font_loaded = True
                                    break
                            if not font_loaded:
                                font = ImageFont.load_default()
                    except Exception as e:
                        logger.error(f"Shrift o'lchamini kichraytirishda xato: {str(e)}")
                        font = ImageFont.load_default()
                    
                    lines = wrap_text(text, font, shape_width)
                    line_height = draw.textbbox((0, 0), "Ay", font=font)[3] - draw.textbbox((0, 0), "Ay", font=font)[1]
                    total_height = len(lines) * line_height
                    logger.info(f"Shrift o'lchami kichraytirildi: {font_size} pt, qatorlar soni: {len(lines)}")

                # Qatorlarni markazlashtirib joylashtirish (2x masshtab)
                top = (shape.top.pt * 2) + (shape_height - total_height) / 2  # Vertikal markazlashtirish
                for line in lines:
                    text_bbox = draw.textbbox((0, 0), line, font=font)
                    text_width = text_bbox[2] - text_bbox[0]
                    left = (shape.left.pt * 2) + (shape_width - text_width) / 2  # Gorizontal markazlashtirish
                    draw.text((left, top), line, font=font, fill=(0, 0, 0))
                    logger.info(f"Qator joylashtirildi: '{line}', joylashuv: ({left}, {top}), shrift: {font_name}, o'lchami: {font_size // 2}")
                    top += line_height

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                logger.debug(f"Processing image shape at ({shape.left.pt}, {shape.top.pt})")
                try:
                    image_stream = shape.image.blob
                    image = Image.open(BytesIO(image_stream))
                    image = image.convert('RGB')
                    img_width, img_height = image.size
                    if shape.width.pt and shape.height.pt:
                        target_width = int(shape.width.pt * 2)
                        target_height = int(shape.height.pt * 2)
                        image = image.resize((target_width, target_height), Image.Resampling.LANCZOS)
                        logger.info(f"Rasm o'lchami moslashtirildi: {target_width}x{target_height}")
                    img.paste(image, (int(shape.left.pt * 2), int(shape.top.pt * 2)))
                    logger.info(f"Rasm joylashtirildi: ({shape.left.pt}, {shape.top.pt})")
                except Exception as e:
                    logger.error(f"Rasmni qayta ishlashda xato: {str(e)}", exc_info=True)

        # QR kod faylini o'chirish
        if qr_code_path and os.path.exists(qr_code_path):
            try:
                os.remove(qr_code_path)
                logger.info(f"QR kod fayli o'chirildi: {qr_code_path}")
            except Exception as e:
                logger.error(f"QR kod faylini o'chirishda xato: {str(e)}")

        img.save(output_path, 'JPEG', quality=100, optimize=False, dpi=(300, 300))
        logger.info(f"Rasm saqlandi: {output_path}")
        return True

    except Exception as e:
        logger.error(f"PPTX to JPG conversion failed: {str(e)}", exc_info=True)
        # QR kod faylini o'chirish (xatolik yuz bersa ham)
        if qr_code_path and os.path.exists(qr_code_path):
            try:
                os.remove(qr_code_path)
                logger.info(f"QR kod fayli o'chirildi: {qr_code_path}")
            except Exception as e:
                logger.error(f"QR kod faylini o'chirishda xato: {str(e)}")
        return False

async def convert_pptx_to_pdf(pptx_path, output_path, context):
    """PPTX dan PDF ga konvertatsiya qilish"""
    try:
        logger.info(f"Starting PPTX to PDF conversion: {pptx_path} -> {output_path}")
        # Avval JPG ga o'tkazamiz
        jpg_path = f"{output_path}.jpg"
        success = await convert_pptx_to_jpg(pptx_path, jpg_path, context)
        if not success:
            logger.error(f"Failed to convert PPTX to JPG: {jpg_path}")
            return False

        # JPG ni PDF ga aylantiramiz (yuqori sifat bilan)
        with open(jpg_path, "rb") as f:
            img_data = f.read()

        with open(output_path, "wb") as f:
            # Yuqori sifatli PDF yaratish
            f.write(img2pdf.convert(img_data, dpi=300))

        logger.info(f"Successfully converted to PDF: {output_path}")
        return True
    except Exception as e:
        logger.error(f"PPTX to PDF conversion error: {str(e)}", exc_info=True)
        return False
    finally:
        if os.path.exists(jpg_path):
            try:
                os.remove(jpg_path)
                logger.info(f"Removed temporary JPG file: {jpg_path}")
            except Exception as e:
                logger.error(f"Error removing JPG file {jpg_path}: {str(e)}")

async def replace_text_and_font(prs, replacements, taqdirlanuvchi_font=None, other_font=None, info_text_font=None, 
                                manzil_va_ega_font=None, shablon_nomi_font=None, shablon_matni_font=None, 
                                shablon_sana_font=None, sana_font=None, diplom_matni_font=None, qr_code_path=None):
    try:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for tag, value in replacements.items():
                        if tag in shape.text:
                            if tag == "{qr_code}":
                                shape.text = shape.text.replace(tag, value)
                                if qr_code_path:
                                    slide.shapes.add_picture(qr_code_path, shape.left, shape.top, width=Pt(100), height=Pt(100))
                                    shape.text = ""
                            else:
                                shape.text = shape.text.replace(tag, value)
                                if tag == "{taqdirlangan}" and taqdirlanuvchi_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = taqdirlanuvchi_font
                                            run.font.size = Pt(24)
                                elif tag == "{sertifikat_matni}" and other_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = other_font
                                            run.font.size = Pt(16)
                                elif tag == "{taqdirlovchi}" and other_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = other_font
                                            run.font.size = Pt(13)
                                elif tag == "{sana}" and sana_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = sana_font
                                            run.font.size = Pt(15)
                                elif tag == "{diplom_matni}" and diplom_matni_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = diplom_matni_font
                                            run.font.size = Pt(16)
                                elif tag == "{info_text}" and info_text_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = info_text_font
                                            run.font.size = Pt(26)
                                elif tag == "{manzil_va_ega}" and manzil_va_ega_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = manzil_va_ega_font
                                            run.font.size = Pt(22)
                                elif tag == "{shablon_nomi}" and shablon_nomi_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = shablon_nomi_font
                                            run.font.size = Pt(26)
                                elif tag == "{shablon_matni}" and shablon_matni_font:
                                    for paragraph in shape.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.name = shablon_matni_font
                                            run.font.size = Pt(18)
                if shape.shape_type == 6:  # Guruh shakli
                    for subshape in shape.shapes:
                        if hasattr(subshape, "text"):
                            for tag, value in replacements.items():
                                if tag in subshape.text:
                                    if tag == "{qr_code}":
                                        subshape.text = subshape.text.replace(tag, value)
                                        if qr_code_path:
                                            slide.shapes.add_picture(qr_code_path, subshape.left, subshape.top, width=Pt(100), height=Pt(100))
                                            subshape.text = ""
                                    else:
                                        subshape.text = subshape.text.replace(tag, value)
                                        if tag == "{taqdirlangan}" and taqdirlanuvchi_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = taqdirlanuvchi_font
                                                    run.font.size = Pt(24)
                                        elif tag == "{sertifikat_matni}" and other_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = other_font
                                                    run.font.size = Pt(16)
                                        elif tag == "{taqdirlovchi}" and other_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = other_font
                                                    run.font.size = Pt(13)
                                        elif tag == "{sana}" and sana_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = sana_font
                                                    run.font.size = Pt(15)
                                        elif tag == "{diplom_matni}" and diplom_matni_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = diplom_matni_font
                                                    run.font.size = Pt(16)
                                        elif tag == "{info_text}" and info_text_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = info_text_font
                                                    run.font.size = Pt(26)
                                        elif tag == "{manzil_va_ega}" and manzil_va_ega_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = manzil_va_ega_font
                                                    run.font.size = Pt(22)
                                        elif tag == "{shablon_nomi}" and shablon_nomi_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = shablon_nomi_font
                                                    run.font.size = Pt(26)
                                        elif tag == "{shablon_matni}" and shablon_matni_font:
                                            for paragraph in subshape.text_frame.paragraphs:
                                                paragraph.alignment = PP_ALIGN.CENTER
                                                for run in paragraph.runs:
                                                    run.font.name = shablon_matni_font
                                                    run.font.size = Pt(18)
        logger.info("Text and font replacement completed successfully")
        return True
    except Exception as e:
        logger.error(f"Error in replace_text_and_font: {str(e)}", exc_info=True)
        raise

async def check_blocked_user(user_id, update, context):
    user_data = load_user_data(user_id)
    if user_data.get('is_blocked', 0) == 1:

        return True
    return False

async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = str(update.effective_user.id)
    if await check_blocked_user(user_id, update, context):
        return ConversationHandler.END
    config = load_config()
    user_data = load_user_data(user_id)

    if not user_data.get('username'):
        user_data = {
            'username': update.effective_user.username or "Nomaâ€˜lum",
            'balance': 5000,
            'referrals': 0,
            'referrals_ids': [],
            'taqdirlangan': "Alisherov Farid",
            'taqdirlovchi': "Bahodirov Dilxush",
            'sertifikat_matni': (
                "Siz bizning Telegram Bot 3-oylik kursimizni 1 - o'rinda muvaffaqiyatli tamomlaganingiz "
                "uchun ushbu sertifikat bilan taqdirlanasiz va sizga jamoamimiz nomidan lutfan tashakkur bildiramiz!"
            ),
            'info_text': (
                "Hurmatli AZIZ MEHMONIMIZ\n"
                "Siz va oila a'zolaringizni 2025-yil 16-fevral kuni soat 18:00 da aziz farzandimiz\n"
                "FARIDJON va SHUKRONAXON larning\n"
                "Nikoh to'yi munosabai bilan yoziladigan dasturxonimizga lutgan taklif etamiz!"
            ),
            'manzil_va_ega': (
                "Hurmat bilan Bahodirovlar oilasi\n"
                "Manzil: Dehqonobod tumani Mamat ota to'yxonasi"
            ),
            'shablon_nomi': "TASHAKKURNOMA",
            'shablon_matni': (
                "Qashqadaryo viloiyati Dehqonobod tumani 87â€“sonli umumta'lim maktabining 5-'A' sinf oâ€˜quvchisi "
                "Alisherov Farid 2024 â€“ 2025 O'quv yilida a'lo va yaxshi baholari hamda namunali xulqi uchun "
                "MAQTOV YORLIG'I bilan taqdirlanadi"
            ),
            'diplom_matni': "Siz bizning kursimizni muvaffaqiyatli tamomlaganingiz uchun ushbu diplom bilan taqdirlanasiz!"
        }
        save_user_data(user_id, user_data)

    keyboard = [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], 
                ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                ["ðŸ“ž Biz bilan aloqa"]]
    await update.message.reply_text(
        "ðŸ‘‹Assalomu alaykum! \n\n"
        "ðŸ¤–Bu bot yordamida siz oson va tez \n ðŸ§¾Sertifikat \n ðŸ’ŒTaklifnoma \nðŸ“‹Shablon \nðŸŽ“Diplom \n Asosida hujjat tayyorlashingiz mumkinâœ….\n"
        "ðŸ“²Jarayonni boshlash uchun quyidagi tugmani bosing:",
        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    )

    if update.message.text.startswith("/start start="):
        referrer_id = update.message.text.split("start=")[1]
        if referrer_id != user_id:
            referrer_data = load_user_data(referrer_id)
            if user_id not in referrer_data.get('referrals_ids', []):
                referrer_data['balance'] = referrer_data.get('balance', 5000) + config['REFERRAL_BONUS']
                referrer_data['referrals'] = referrer_data.get('referrals', 0) + 1
                referrer_data['referrals_ids'] = referrer_data.get('referrals_ids', []) + [user_id]
                save_user_data(referrer_id, referrer_data)
                await update.message.reply_text(
                    f"âœ… Referal orqali roâ€˜yxatdan oâ€˜tdingiz! Referal egasiga {config['REFERRAL_BONUS']} soâ€˜m qoâ€˜shildi."
                )
            else:
                await update.message.reply_text(
                    "â„¹ï¸ Siz allaqachon ushbu referal linki orqali roâ€˜yxatdan oâ€˜tgansiz."
                )

    return START

async def handle_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if await check_blocked_user(user_id, update, context):
        return ConversationHandler.END
    config = load_config()
    if update.message.text.startswith("ðŸŸ¢ Sertifikat tayyorlash"):
        user_data = load_user_data(user_id)
        if user_data.get('balance', 5000) < config['CERTIFICATE_COST']:
            await update.message.reply_text(
                f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q. Joriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
                f"ðŸ’µSertifikat narxi: {config['CERTIFICATE_COST']} soâ€˜m.\nðŸ’µIltimos, balansingizni toâ€˜ldiring.",
                reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
            )            
            return BALANCE
        templates = get_templates(SERTIFIKAT_PATH)
        if not os.path.exists(SAMPLE_IMAGE):
            await update.message.reply_text("âŒ namuna.jpg fayli topilmadi.")
            return ConversationHandler.END
        keyboard = create_template_keyboard(templates)
        if not keyboard:
            await update.message.reply_text("âŒ Sertifikat papkasida mos .pptx fayllari topilmadi.")
            return ConversationHandler.END
        with open(SAMPLE_IMAGE, "rb") as img, open(SAMPLE_IMAGE1, "rb") as img1, open(SAMPLE_IMAGE2, "rb") as img2, open(SAMPLE_IMAGE3, "rb") as img3, open(SAMPLE_IMAGE4, "rb") as img4:
            await update.message.reply_photo(photo=img)
            await update.message.reply_photo(photo=img1)
            await update.message.reply_photo(photo=img2)
            await update.message.reply_photo(photo=img3)
            await update.message.reply_photo(
                photo=img4,
                caption="Quyidagi namunalardan birini tanlang:",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        context.user_data['templates'] = templates
        context.user_data['doc_type'] = 'Sertifikat'
        return TEMPLATE
    elif update.message.text.startswith("ðŸŸ¢ Taklifnoma yasash"):
        user_data = load_user_data(user_id)
        if user_data.get('balance', 5000) < config['TAKLIFNOMA_COST']:
            await update.message.reply_text(
                f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q. Joriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
                f"ðŸ’µTaklifnoma narxi: {config['TAKLIFNOMA_COST']} soâ€˜m.\nðŸ’µIltimos, balansingizni toâ€˜ldiring.",
                reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
            )
            return BALANCE
        if not os.path.exists(TAKLIFNOMA_SAMPLE_IMAGE):
            await update.message.reply_text("âŒ taklifnoma_namuna.jpg fayli topilmadi.")
            return ConversationHandler.END
        keyboard = create_template_keyboard()
        with open(TAKLIFNOMA_SAMPLE_IMAGE, "rb") as img, open(TAKLIFNOMA_SAMPLE_IMAGE1, "rb") as img1, open(TAKLIFNOMA_SAMPLE_IMAGE2, "rb") as img2:
            await update.message.reply_photo(photo=img)
            await update.message.reply_photo(photo=img1)
            await update.message.reply_photo(
                photo=img2,
                caption="Quyidagi namunalardan birini tanlang (1-10):",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        context.user_data['doc_type'] = 'Taklifnoma'
        return TEMPLATE
    elif update.message.text.startswith("ðŸŸ¢ Namuna orqali yasash"):
        user_data = load_user_data(user_id)
        if user_data.get('balance', 5000) < config['SHABLON_COST']:
            await update.message.reply_text(
                f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q. Joriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
                f"ðŸ’µShablon narxi: {config['SHABLON_COST']} soâ€˜m.\nðŸ’°Iltimos, balansingizni toâ€˜ldiring.",
                reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
            )
            return BALANCE
        if not os.path.exists(SHABLON_IMAGE):
            await update.message.reply_text("âŒ shablon.jpg fayli topilmadi.")
            return ConversationHandler.END
        keyboard = create_template_keyboard()
        with open(SHABLON_IMAGE, "rb") as img, open(SHABLON_IMAGE1, "rb") as img1, open(SHABLON_IMAGE2, "rb") as img2:
            await update.message.reply_photo(photo=img)
            await update.message.reply_photo(photo=img1)
            await update.message.reply_photo(
                photo=img2,
                caption="Quyidagi namunalardan birini tanlang (1-15):",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        context.user_data['doc_type'] = 'Shablon'
        return SHABLON
    elif update.message.text.startswith("ðŸŸ¢ Diplom yasash"):
        user_data = load_user_data(user_id)
        if user_data.get('balance', 5000) < config['DIPLOM_COST']:
            await update.message.reply_text(
                f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q. Joriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
                f"ðŸ’µDiplom narxi: {config['DIPLOM_COST']} soâ€˜m.\nðŸ’µIltimos, balansingizni toâ€˜ldiring.",
                reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
            )
            return BALANCE
        templates = get_templates(DIPLOM_PATH)
        if not os.path.exists(DIPLOM_NAMUNA_IMAGE):
            await update.message.reply_text("âŒ diplom_namuna.jpg fayli topilmadi.")
            return ConversationHandler.END
        keyboard = create_template_keyboard(templates)
        if not keyboard:
            await update.message.reply_text("âŒ Diplom papkasida mos .pptx fayllari topilmadi.")
            return ConversationHandler.END
        with open(DIPLOM_NAMUNA_IMAGE, "rb") as img, open(DIPLOM_NAMUNA_IMAGE1, "rb") as img1, open(DIPLOM_NAMUNA_IMAGE2, "rb") as img2, open(DIPLOM_NAMUNA_IMAGE3, "rb") as img3, open(DIPLOM_NAMUNA_IMAGE4, "rb") as img4:
            await update.message.reply_photo(photo=img)
            await update.message.reply_photo(photo=img1)
            await update.message.reply_photo(photo=img2)
            await update.message.reply_photo(photo=img3)
            await update.message.reply_photo(
                photo=img4,
                caption="Quyidagi namunalardan birini tanlang (1-26):",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        context.user_data['templates'] = templates
        context.user_data['doc_type'] = 'Diplom'
        return TEMPLATE
    elif update.message.text.startswith("ðŸ’°"):
        user_data = load_user_data(user_id)
        await update.message.reply_text(
            f"ðŸ’° Joriy balansingiz: {user_data.get('balance', 5000)} soâ€˜m\n"
            f"ðŸ–‡Chaqirilgan referallar: {user_data.get('referrals', 0)}\n"
            f"ðŸ“ŠUmumiy natijalaringiz soni: {user_data.get('result_count', 0)}\n"  
            f"ðŸ§¾Sertifikat narxi: {config['CERTIFICATE_COST']} soâ€˜m\n"
            f"ðŸ’ŒTaklifnoma narxi: {config['TAKLIFNOMA_COST']} soâ€˜m\n"
            f"ðŸ“œNamuna orqali yasash narxi: {config['SHABLON_COST']} soâ€˜m\n"
            f"ðŸŽ“Diplom narxi: {config['DIPLOM_COST']} soâ€˜m\n"
            f"ðŸ—‚PDF qilish narxi: {config['PDF_COST']} soâ€˜m\n"
            f"ðŸ”—Referral uchun Bonus: {config['REFERRAL_BONUS']} soâ€˜m\n \n"
            f"ðŸ”—Referal link: https://t.me/sertifikat_tayyorlashbot?start=start={user_id}",
            reply_markup=ReplyKeyboardMarkup([["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], 
                                             ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                                             ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                                             ["ðŸ“ž Biz bilan aloqa"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return START
    elif update.message.text.startswith("ðŸ’¸"):
        keyboard = [["ðŸ’³ Kartaga toâ€˜lov", "ðŸ¤ Referal orqali"], ["â¬…ï¸ Orqaga"]]
        await update.message.reply_text(
            "ðŸ’°Balansni toâ€˜ldirish usulini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return PAYMENT_METHOD
    elif update.message.text.startswith("ðŸ“ž Biz bilan aloqa"):
        await update.message.reply_text(
            "ðŸ“ž Murojaatingizni yuboring (matn yoki rasm):",
            reply_markup=ReplyKeyboardRemove()
        )
        return CONTACT
    elif context.user_data.get('awaiting_topup') and user_id == config['ADMIN_ID']:
        return await handle_admin_topup(update, context)
    else:
        keyboard = [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], 
                    ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                    ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                    ["ðŸ“ž Biz bilan aloqa"]]
        await update.message.reply_text(
            "ðŸ“²Iltimos, quyidagi tugmalardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return START

async def handle_contact(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()

    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text or update.message.photo:
        return await handle_contact_message(update, context)
    await update.message.reply_text(
        "ðŸ“ž Murojaatingizni matn sifatida yuboring:",
        reply_markup=ReplyKeyboardRemove()
    )
    return CONTACT_MESSAGE


async def handle_contact_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    user_data = load_user_data(user_id)
    username = user_data.get('username', 'Nomaâ€˜lum')

    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)

    try:
        if update.message.text:
            await context.bot.send_message(
                chat_id=config['ADMIN_ID'],
                text=f"Murojaat\nFoydalanuvchi ID: {user_id}\nUsername: {username}\nXabar: {update.message.text}"
            )
            await update.message.reply_text(
                "âœ… Xabaringiz adminga yetkazildi!",
                reply_markup=ReplyKeyboardMarkup(
                    [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                     ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                     ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                     ["ðŸ“ž Biz bilan aloqa"]],
                    one_time_keyboard=True, resize_keyboard=True
                )
            )
            return START
        else:
            await update.message.reply_text(
                "âŒ Faqat matn xabari yuborishingiz mumkin. Iltimos, matn sifatida murojaatingizni yuboring.",
                reply_markup=ReplyKeyboardRemove()
            )
            return CONTACT_MESSAGE

    except Exception as e:
        await update.message.reply_text(
            f"âŒ Xabarni adminga yuborishda xatolik yuz berdi: {str(e)}",
            reply_markup=ReplyKeyboardMarkup(
                [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                 ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                 ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                 ["ðŸ“ž Biz bilan aloqa"]],
                one_time_keyboard=True, resize_keyboard=True
            )
        )
        return START

def update_balance(user_id: str, amount: int, transaction_type: str) -> bool:
    """Foydalanuvchi balansini sinxron ravishda yangilaydi."""
    with db_lock:
        try:
            user_data = load_user_data(user_id)
            new_balance = user_data['balance'] + amount
            if new_balance < 0:
                logger.warning(f"Balans yetarli emas: user_id={user_id}, amount={amount}")
                return False
            user_data['balance'] = new_balance
            save_user_data(user_id, user_data)
            log_transaction(user_id, amount, transaction_type)
            logger.info(f"Balans yangilandi: user_id={user_id}, new_balance={new_balance}, type={transaction_type}")
            return True
        except Exception as e:
            logger.error(f"Balansni yangilashda xato: user_id={user_id}, {str(e)}")
            raise
        
async def handle_balance(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text.startswith("ðŸ’¸"):
        keyboard = [["ðŸ’³ Kartaga toâ€˜lov", "ðŸ¤ Referal orqali"], ["â¬…ï¸ Orqaga"]]
        await update.message.reply_text(
            "Balansni toâ€˜ldirish usulini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return PAYMENT_METHOD
    elif update.message.text.startswith("ðŸŸ¢"):
        return await handle_start(update, context)
    else:
        keyboard = [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], ["ðŸŸ¢ Namuna orqali yasash", "ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"]]
        await update.message.reply_text(
            "Iltimos, quyidagi tugmalardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return START

async def handle_payment_method(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if update.message.text.startswith("ðŸ’³"):
        await update.message.reply_text(
            f"ðŸ’³ Karta raqami: <code>{config['CARD_NUMBER']}</code>\n"
            f"ðŸ“Œ Karta egasi: {config.get('CARD_HOLDER', 'Nomaâ€˜lum')}\n"
            f"Iltimos, xohlagan summangizni ushbu kartaga oâ€˜tkazing va <code>âœ… Toâ€˜lov qildim</code> tugmani bosing:",
            parse_mode="HTML",
            reply_markup=ReplyKeyboardMarkup([["âœ… Toâ€˜lov qildim"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return CARD_PAYMENT
    elif update.message.text.startswith("ðŸ¤"):
        user_data = load_user_data(user_id)
        await update.message.reply_text(
            f"ðŸ”— Sizning referal linkingiz:\n \n"
            f"ðŸ“Œ https://t.me/sertifikat_tayyorlashbot?start=start={user_id}\n \n"
            f"ðŸ–‡ Chaqirilgan referallar: {user_data.get('referrals', 0)}\n"
            "ðŸ”— Ushbu linkni doâ€˜stlaringizga yuboring. Har bir yangi foydalanuvchi uchun balansingizga to'ldiring!",
            reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return PAYMENT_METHOD
    else:
        user_data = load_user_data(user_id)
        await update.message.reply_text(
            f"ðŸ’° Joriy balansingiz: {user_data.get('balance', 5000)} soâ€˜m\n"
            f"ðŸ–‡ Chaqirilgan referallar: {user_data.get('referrals', 0)}\n"
            f"ðŸ§¾ Sertifikat narxi: {config['CERTIFICATE_COST']} soâ€˜m\n"
            f"ðŸ’Œ Taklifnoma narxi: {config['TAKLIFNOMA_COST']} soâ€˜m\n"
            f"ðŸ“œ Namuna orqali yasash narxi: {config['SHABLON_COST']} soâ€˜m\n"
            f"ðŸ—‚ PDF qilish narxi: {config['PDF_COST']} soâ€˜m",
            reply_markup=ReplyKeyboardMarkup([["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], ["ðŸŸ¢ Namuna orqali yasash", "ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return START

async def handle_card_payment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text.startswith("âœ…"):
        await update.message.reply_text(
            "Iltimos, toâ€˜lov chekini Rasm ðŸ–¼(JPG,PNG,Screenshoot) yoki PDF ðŸ—‚ formatida yuboring ðŸ“¥.",
            reply_markup=ReplyKeyboardRemove()
        )
        return UPLOAD_CHECK
    return PAYMENT_METHOD

async def handle_upload_check(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.document or update.message.photo:
        if update.message.document:
            file = update.message.document
            if file.mime_type not in ["image/jpeg", "image/png", "application/pdf"]:
                await update.message.reply_text(
                    "âŒ Faqat ðŸ–¼Rasm (JPG, PNG, Skrinshot) yoki ðŸ—‚PDF formatidagi fayllar qabul qilinadiðŸ“¥. Iltimos, qaytadan yuboringðŸ“¤."
                )
                return UPLOAD_CHECK
            file_id = file.file_id
            context.user_data['check_file_type'] = 'document'
        elif update.message.photo:
            file_id = update.message.photo[-1].file_id
            context.user_data['check_file_type'] = 'photo'
        try:
            context.user_data['check_file_id'] = file_id
            await update.message.reply_text(
                "âœ… Chek qabul qilindiðŸ“¤. Iltimos, toâ€˜lov haqida qisqacha kommentariya yozingðŸ“ (masalan, toâ€˜langan summa):",
                reply_markup=ReplyKeyboardRemove()
            )
            return COMMENT
        except Exception as e:
            await update.message.reply_text(
                "âŒ Chekni qayta ishlashda xatolik yuz berdi. Iltimos, qaytadan yuboring."
            )
            return UPLOAD_CHECK
    else:
        await update.message.reply_text(
            "âŒ Iltimos, toâ€˜lov chekini Rasm (JPG, PNG, Skrinshot) yoki fayl (PDF) sifatida yuboring."
        )
        return UPLOAD_CHECK

async def handle_comment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    comment = update.message.text.strip()
    check_file_id = context.user_data.get('check_file_id')
    check_file_type = context.user_data.get('check_file_type')
    if not check_file_id or not check_file_type:
        await update.message.reply_text(
            "âŒ Chek fayli topilmadi. Iltimos, chekni qaytadan yuboring.",
            reply_markup=ReplyKeyboardMarkup([["âœ… Toâ€˜lov qildim"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return CARD_PAYMENT
    try:
        user_data = load_user_data(user_id)
        username = user_data.get('username', 'Nomaâ€˜lum')
        keyboard = [[InlineKeyboardButton("ðŸ’¸ Balansni toâ€˜ldirish", callback_data=f"topup_{user_id}")]]
        reply_markup = InlineKeyboardMarkup(keyboard)
        context.user_data['topup_user_id'] = user_id
        caption = f"Toâ€˜lov cheki\nFoydalanuvchi ID: {user_id}\nUsername: {username}\nKomentariya: {comment}"
        if check_file_type == 'document':
            await context.bot.send_document(
                chat_id=config['ADMIN_ID'],
                document=check_file_id,
                caption=caption,
                reply_markup=reply_markup
            )
        else:
            await context.bot.send_photo(
                chat_id=config['ADMIN_ID'],
                photo=check_file_id,
                caption=caption,
                reply_markup=reply_markup
            )
        await update.message.reply_text(
            "âœ… Toâ€˜lov cheki va kommentariya adminga yuborildiðŸ“§. Tasdiqlashni kutingðŸ•˜.",
            reply_markup=ReplyKeyboardMarkup([["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"], ["ðŸŸ¢ Namuna orqali yasash", "ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return START
    except Exception as e:
        await update.message.reply_text(
            "âŒ Chekni adminga yuborishda xatolik yuz berdi. Iltimos, qaytadan urinib koâ€˜ring yoki Adminga murojaat qiling",
            reply_markup=ReplyKeyboardMarkup([["âœ… Toâ€˜lov qildim"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return CARD_PAYMENT

async def handle_callback_query(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await query.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return
    if query.data.startswith("topup_"):
        target_user_id = query.data.split("_")[1]
        context.user_data['topup_user_id'] = target_user_id
        context.user_data['awaiting_topup'] = True
        await query.message.reply_text(
            f"Foydalanuvchi ID {target_user_id} uchun qancha summa qoâ€˜shilsin? (soâ€˜mda kiriting)",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_TOPUP

async def handle_admin_topup(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END
    try:
        amount = int(update.message.text.strip())
        target_user_id = context.user_data.get('topup_user_id')
        if not target_user_id:
            await update.message.reply_text(
                "âŒ Foydalanuvchi ID topilmadi. Iltimos, chek xabaridan tugmani bosing.",
                reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
            )
            return ADMIN_PANEL
        user_data = load_user_data(target_user_id)
        user_data['balance'] = user_data.get('balance', 5000) + amount
        save_user_data(target_user_id, user_data)
        context.user_data['awaiting_topup'] = False
        await update.message.reply_text(
            f"âœ… ID {target_user_id} balansiga {amount} soâ€˜m qoâ€˜shildi.",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        await context.bot.send_message(
            chat_id=target_user_id,
            text=f"ðŸ’° Balansingiz {amount} soâ€˜mga toâ€˜ldirildi.\n Joriy balans: {user_data['balance']} soâ€˜m"
        )
        return ADMIN_PANEL
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return ADMIN_TOPUP

async def admin_panel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    keyboard = [
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ]
    await update.message.reply_text(
        "Admin panel:\nQuyidagi amallardan birini tanlang:",
        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    )
    return ADMIN_PANEL

async def handle_admin_panel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    text = update.message.text
    if text == "ðŸ“Š Statistika":
        with get_db_connection() as conn:
            c = conn.cursor()
            c.execute("SELECT COUNT(*) as user_count FROM users")
            user_count = c.fetchone()['user_count']
            c.execute("SELECT COUNT(*) as blocked_count FROM users WHERE is_blocked = 1")
            blocked_count = c.fetchone()['blocked_count']
            c.execute("SELECT COUNT(*), SUM(amount) FROM transactions")
            transaction_data = c.fetchone()
            transaction_count = transaction_data[0]
            total_amount = transaction_data[1] or 0
            c.execute("SELECT COUNT(*) as referral_count FROM users WHERE referrals > 0")
            referral_count = c.fetchone()['referral_count']

        stats_text = (
            f"ðŸ“Š Statistika:\n"
            f"ðŸ‘¥ Foydalanuvchilar soni: {user_count}\n"
            f"ðŸš« Bloklangan foydalanuvchilar: {blocked_count}\n"
            f"ðŸ’¸ Tranzaksiyalar soni: {transaction_count}\n"
            f"ðŸ’° Umumiy tranzaksiya summasi: {total_amount} soâ€˜m\n"
            f"ðŸ–‡ Referallar soni: {referral_count}"
        )
        await update.message.reply_text(
            stats_text,
            reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL
    elif text == "ðŸ‘¤ Foydalanuvchi":
        await update.message.reply_text(
            "Foydalanuvchi ID yoki username kiriting:",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_FOYDALANUVCHI
    elif update.message.text.startswith("âš™ï¸ Narxlarni sozlash"):
        keyboard = [
            [f"Admin - {config['ADMIN_ID']}", f"Karta raqami - {config['CARD_NUMBER']}"],
            [f"Karta egasi - {config.get('CARD_HOLDER', 'Nomalum')}"],
            [f"Shablon narxi - {config['SHABLON_COST']}", f"Sertifikat narxi - {config['CERTIFICATE_COST']}"],
            [f"Taklifnoma narxi - {config['TAKLIFNOMA_COST']}", f"Diplom narxi - {config['DIPLOM_COST']}"],
            [f"PDF qilish narxi - {config['PDF_COST']}", f"Referal narxi - {config['REFERRAL_BONUS']}"],
            ["â¬…ï¸ Orqaga"]
        ]
        await update.message.reply_text(
            "âš™ï¸ Narxlarni sozlash boâ€˜limi:\nQuyidagi parametrlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return CONFIG_PRICE
    elif update.message.text.startswith("ðŸ“‹"):
        with get_db_connection() as conn:
            c = conn.cursor()
            c.execute("SELECT user_id, username, balance, result_count FROM users")
            users = c.fetchall()

        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerow(["User ID", "Username", "Balance", "Count"])  
        for user in users:
            writer.writerow(user)
        output.seek(0)
        await context.bot.send_document(
            chat_id=user_id,
            document=io.BytesIO(output.getvalue().encode('utf-8')),
            filename="users.csv",
            caption="Foydalanuvchilar roâ€˜yxati"
        )
        return ADMIN_PANEL
    elif update.message.text.startswith("ðŸ’¸"):
        await update.message.reply_text(
            "Balansni toâ€˜ldirish uchun foydalanuvchi username yoki user ID kiriting:",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_USER
    elif update.message.text.startswith("ðŸ“¬"):
        keyboard = [["ðŸ“¢ Barcha foydalanuvchilarga", "ðŸ‘¤ Alohida foydalanuvchiga"], ["â¬…ï¸ Orqaga"]]
        await update.message.reply_text(
            "ðŸ“¬ Xabar yuborish turini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_MESSAGE_TYPE
    elif update.message.text == "/start":
        return await start(update, context)
    else:
        keyboard = [
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ]
        await update.message.reply_text(
            "Iltimos, quyidagi tugmalardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

async def handle_config_price(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END
    if update.message.text == "â¬…ï¸ Orqaga":
        return await admin_panel(update, context)
    text = update.message.text.strip()
    if text.startswith("Karta egasi -"):
        await update.message.reply_text(
            "Yangi karta egasi ismini kiriting (masalan, Bahodirov D):",
            reply_markup=ReplyKeyboardRemove()
        )
        context.user_data['config_key'] = 'CARD_HOLDER'
        return SET_NEW_PRICE
    selected_option = update.message.text.strip()
    price_keys = {
        f"Admin - {config['ADMIN_ID']}": "ADMIN_ID",
        f"Karta raqami - {config['CARD_NUMBER']}": "CARD_NUMBER",
        f"Shablon narxi - {config['SHABLON_COST']}": "SHABLON_COST",
        f"Sertifikat narxi - {config['CERTIFICATE_COST']}": "CERTIFICATE_COST",
        f"Taklifnoma narxi - {config['TAKLIFNOMA_COST']}": "TAKLIFNOMA_COST",
        f"Diplom narxi - {config['DIPLOM_COST']}": "DIPLOM_COST",
        f"PDF qilish narxi - {config['PDF_COST']}": "PDF_COST",
        f"Referal narxi - {config['REFERRAL_BONUS']}": "REFERRAL_BONUS"
    }
    if selected_option in price_keys:
        context.user_data['selected_price_key'] = price_keys[selected_option]
        await update.message.reply_text(
            f"Yangi qiymatni kiriting (raqamda):",
            reply_markup=ReplyKeyboardRemove()
        )
        return SET_NEW_PRICE
    await update.message.reply_text(
        "Iltimos, quyidagi parametrlardan birini tanlang:",
        reply_markup=ReplyKeyboardMarkup(
            [
                [f"Admin - {config['ADMIN_ID']}", f"Karta raqami - {config['CARD_NUMBER']}"],
                [f"Karta egasi - {config['CARD_HOLDER']}"],
                [f"Shablon narxi - {config['SHABLON_COST']}", f"Sertifikat narxi - {config['CERTIFICATE_COST']}"],
                [f"Taklifnoma narxi - {config['TAKLIFNOMA_COST']}", f"Diplom narxi - {config['DIPLOM_COST']}"],
                [f"PDF qilish narxi - {config['PDF_COST']}", f"Referal narxi - {config['REFERRAL_BONUS']}"],
                ["â¬…ï¸ Orqaga"]
            ],
            one_time_keyboard=True, resize_keyboard=True
        )
    )
    return CONFIG_PRICE

async def handle_set_new_price(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END
    text = update.message.text.strip()
    config_key = context.user_data.get('config_key')
    try:
        if config_key == 'CARD_HOLDER':
            config['CARD_HOLDER'] = text
            save_config(config)
            await update.message.reply_text(
                f"âœ… Karta egasi {text} ga oâ€˜zgartirildi.",
                reply_markup=ReplyKeyboardMarkup(
                    [
                        [f"Admin - {config['ADMIN_ID']}", f"Karta raqami - {config['CARD_NUMBER']}"],
                        [f"Karta egasi - {config['CARD_HOLDER']}"],
                        [f"Shablon narxi - {config['SHABLON_COST']}", f"Sertifikat narxi - {config['CERTIFICATE_COST']}"],
                        [f"Taklifnoma narxi - {config['TAKLIFNOMA_COST']}", f"Diplom narxi - {config['DIPLOM_COST']}"],
                        [f"PDF qilish narxi - {config['PDF_COST']}", f"Referal narxi - {config['REFERRAL_BONUS']}"],
                        ["â¬…ï¸ Orqaga"]
                    ],
                    one_time_keyboard=True, resize_keyboard=True
                )
            )
            return CONFIG_PRICE
    except ValueError:
        await update.message.reply_text("Iltimos, toâ€˜gâ€˜ri qiymat kiriting.")
        return SET_NEW_PRICE
    try:
        new_value = update.message.text.strip()
        if not new_value.isdigit():
            await update.message.reply_text("âŒ Iltimos, faqat raqam kiriting.")
            return SET_NEW_PRICE
        selected_key = context.user_data.get('selected_price_key')
        if not selected_key:
            await update.message.reply_text(
                "âŒ Parametr tanlanmadi. Iltimos, qaytadan boshlang.",
                reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
            )
            return ADMIN_PANEL
        config[selected_key] = int(new_value) if selected_key != "CARD_NUMBER" else new_value
        save_config(config)
        await update.message.reply_text(
            f"âœ… {selected_key} uchun yangi qiymat ({new_value}) muvaffaqiyatli saqlandi.",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL
    except Exception as e:
        await update.message.reply_text(
            f"âŒ Xatolik yuz berdi: {str(e)}",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

async def handle_admin_message_type(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END
    if update.message.text == "â¬…ï¸ Orqaga":
        return await admin_panel(update, context)
    elif update.message.text == "ðŸ“¢ Barcha foydalanuvchilarga":
        context.user_data['message_recipient'] = 'all'
        await update.message.reply_text(
            "ðŸ“¬ Yuboriladigan xabarni kiriting (matn, rasm, video yoki forward qilingan xabar):",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_MESSAGE_CONTENT
    elif update.message.text == "ðŸ‘¤ Alohida foydalanuvchiga":
        await update.message.reply_text(
            "Foydalanuvchi username yoki user ID kiriting:",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_MESSAGE_RECIPIENT
    else:
        keyboard = [["ðŸ“¢ Barcha foydalanuvchilarga", "ðŸ‘¤ Alohida foydalanuvchiga"], ["â¬…ï¸ Orqaga"]]
        await update.message.reply_text(
            "Iltimos, quyidagi tugmalardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_MESSAGE_TYPE

async def handle_admin_message_recipient(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END
    identifier = update.message.text.strip()
    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("SELECT user_id FROM users WHERE username = ? OR user_id = ?", (identifier, identifier))
        user = c.fetchone()

    if not user:
        await update.message.reply_text(
            f"âŒ {identifier} topilmadi. Iltimos, toâ€˜gâ€˜ri username yoki user ID kiriting.",
            reply_markup=ReplyKeyboardMarkup([["ðŸ“¢ Barcha foydalanuvchilarga", "ðŸ‘¤ Alohida foydalanuvchiga"], ["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_MESSAGE_TYPE
    context.user_data['message_recipient'] = user[0]
    await update.message.reply_text(
        "ðŸ“¬ Yuboriladigan xabarni kiriting (matn, rasm, video yoki forward qilingan xabar):",
        reply_markup=ReplyKeyboardRemove()
    )
    return ADMIN_MESSAGE_CONTENT

async def handle_admin_message_content(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    recipient = context.user_data.get('message_recipient')
    if not recipient:
        await update.message.reply_text(
            "âŒ Xabar yuborish uchun foydalanuvchi tanlanmadi. Qaytadan boshlang.",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

    try:
        with get_db_connection() as conn:
            c = conn.cursor()
            if recipient == 'all':
                c.execute("SELECT user_id FROM users")
                recipients = [row[0] for row in c.fetchall()]
            else:
                recipients = [recipient]

        for recipient_id in recipients:
            try:
                if update.message.text:
                    await context.bot.send_message(
                        chat_id=recipient_id,
                        text=update.message.text,
                        parse_mode="HTML",
                        reply_markup=update.message.reply_markup if update.message.reply_markup else None
                    )
                elif update.message.photo:
                    await context.bot.send_photo(
                        chat_id=recipient_id,
                        photo=update.message.photo[-1].file_id,
                        caption=update.message.caption,
                        parse_mode="HTML",
                        reply_markup=update.message.reply_markup if update.message.reply_markup else None
                    )
                elif update.message.video:
                    await context.bot.send_video(
                        chat_id=recipient_id,
                        video=update.message.video.file_id,
                        caption=update.message.caption,
                        parse_mode="HTML",
                        reply_markup=update.message.reply_markup if update.message.reply_markup else None
                    )
                elif update.message.forward_from or update.message.forward_from_chat:
                    await context.bot.forward_message(
                        chat_id=recipient_id,
                        from_chat_id=update.message.chat_id,
                        message_id=update.message.message_id
                    )
            except Exception as e:
                await update.message.reply_text(
                    f"âŒ ID {recipient_id} ga xabar yuborishda xatolik: {str(e)}"
                )

        await update.message.reply_text(
            f"âœ… Xabar {len(recipients)} ta foydalanuvchiga yuborildi!",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL
    except Exception as e:
        await update.message.reply_text(
            f"âŒ Xabar yuborishda xatolik yuz berdi: {str(e)}",
            reply_markup=ReplyKeyboardMarkup([
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

async def handle_foydalanuvchi(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    search_query = update.message.text.strip()
    user_data = None

    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("SELECT * FROM users WHERE user_id = ? OR username = ?", (search_query, search_query))
        user = c.fetchone()
        if user:
            user_data = {
                "user_id": user['user_id'],
                "username": user['username'],
                "balance": user['balance'],
                "referrals": user['referrals'],
                "referrals_ids": json.loads(user['referrals_ids']),
                "result_count": user['result_count'],
                "is_blocked": user['is_blocked'],
                "taqdirlangan": user['taqdirlangan'],
                "taqdirlovchi": user['taqdirlovchi'],
                "sertifikat_matni": user['sertifikat_matni'],
                "sana": user['sana'],
                "info_text": user['info_text'],
                "manzil_va_ega": user['manzil_va_ega'],
                "shablon_nomi": user['shablon_nomi'],
                "shablon_matni": user['shablon_matni'],
                "diplom_matni": user['diplom_matni']
            }

    if not user_data:
        await update.message.reply_text(
            "âŒ Foydalanuvchi topilmadi. Iltimos, toâ€˜gâ€˜ri ID yoki username kiriting.",
            reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

    info_text = (
        f"ðŸ‘¤ Foydalanuvchi ma'lumotlari:\n"
        f"ID: {user_data['user_id']}\n"
        f"Username: {user_data['username']}\n"
        f"Balans: {user_data['balance']} soâ€˜m\n"
        f"Referallar soni: {user_data['referrals']}\n"
        f"Natijalar soni: {user_data['result_count']}\n"
        f"Bloklangan: {'Ha' if user_data['is_blocked'] else 'Yoâ€˜q'}\n"
        f"Taqdirlangan: {user_data['taqdirlangan'] or 'Kiritilmagan'}\n"
        f"Taqdirlovchi: {user_data['taqdirlovchi'] or 'Kiritilmagan'}\n"
        f"Sertifikat matni: {user_data['sertifikat_matni'] or 'Kiritilmagan'}\n"
        f"Diplom matni: {user_data['diplom_matni'] or 'Kiritilmagan'}\n"
        f"Taklifnoma matni: {user_data['info_text'] or 'Kiritilmagan'}\n"
        f"Manzil va ega: {user_data['manzil_va_ega'] or 'Kiritilmagan'}\n"
        f"Shablon nomi: {user_data['shablon_nomi'] or 'Kiritilmagan'}\n"
        f"Shablon matni: {user_data['shablon_matni'] or 'Kiritilmagan'}\n"
        f"Sana: {user_data['sana'] or 'Kiritilmagan'}"
    )

    keyboard = [
        ["ðŸ”¼ Balansni toâ€˜ldirish", "ðŸ’¬ Xabar yuborish"],
        [f"{'âŒ Bloklash' if not user_data['is_blocked'] else 'ðŸ”“ Blokdan chiqarish'}"],
        ["â¬…ï¸ Orqaga"]
    ]
    context.user_data['selected_user_id'] = user_data['user_id']
    await update.message.reply_text(
        info_text,
        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    )
    return ADMIN_FOYDALANUVCHI

async def handle_admin_user(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()

    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    identifier = update.message.text.strip()

    with get_db_connection() as conn:
        c = conn.cursor()
        c.execute("SELECT user_id FROM users WHERE user_id = ? OR username = ?", (identifier, identifier))
        result = c.fetchone()

    if not result:
        await update.message.reply_text(
            "âŒ Bunday foydalanuvchi topilmadi. Iltimos, toâ€˜gâ€˜ri user ID yoki username kiriting."
        )
        return ADMIN_USER

    context.user_data['topup_user_id'] = result['user_id']
    await update.message.reply_text(
        f"{identifier} (ID: {result['user_id']}) uchun qancha summa qoâ€˜shilsin? (soâ€˜mda kiriting)",
        reply_markup=ReplyKeyboardRemove()
    )
    return ADMIN_TOPUP

async def handle_admin_action(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    config = load_config()
    if user_id != config['ADMIN_ID']:
        await update.message.reply_text("âŒ Sizda admin huquqlari yoâ€˜q.")
        return ConversationHandler.END

    text = update.message.text
    selected_user_id = context.user_data.get('selected_user_id')

    if not selected_user_id:
        await update.message.reply_text(
            "âŒ Foydalanuvchi tanlanmadi. Iltimos, foydalanuvchi ID yoki username kiriting.",
            reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

    if text == "ðŸ”¼ Balansni toâ€˜ldirish":
        await update.message.reply_text(
            f"Foydalanuvchi ID {selected_user_id} uchun qancha summa qoâ€˜shilsin? (soâ€˜mda kiriting, masalan, 10000):",
            reply_markup=ReplyKeyboardRemove()
        )
        return ADMIN_TOPUP

    elif text == "ðŸ’¬ Xabar yuborish":
        keyboard = [["ðŸ“¢ Umumiy xabar"], [f"ðŸ‘¤ Foydalanuvchiga xabar ({selected_user_id})"], ["â¬…ï¸ Orqaga"]]
        await update.message.reply_text(
            "Xabar turini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_MESSAGE_TYPE

    elif text in ("âŒ Bloklash", "ðŸ”“ Blokdan chiqarish"):
        user_data = load_user_data(selected_user_id)
        is_blocked = 1 if text == "âŒ Bloklash" else 0
        user_data['is_blocked'] = is_blocked
        save_user_data(selected_user_id, user_data)
        await update.message.reply_text(
            f"âœ… Foydalanuvchi {selected_user_id} {'bloklandi' if is_blocked else 'blokdan chiqarildi'}.",
            reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

    elif text == "â¬…ï¸ Orqaga":
        keyboard = [
        ["ðŸ“Š Statistika", "ðŸ“‹ Foydalanuvchilar roâ€˜yxati"],
        ["ðŸ“¬ Xabar yuborish", "ðŸ‘¤ Foydalanuvchi"],
        ["âš™ï¸ Narxlarni sozlash", "ðŸ’¸ Balans toâ€˜ldirish"],
        ["â¬…ï¸ Orqaga"]
    ]
        await update.message.reply_text(
            "Admin panel:\nQuyidagi amallardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return ADMIN_PANEL

    await update.message.reply_text(
        "Notoâ€˜gâ€˜ri buyruq. Iltimos, tugmalardan birini tanlang.",
        reply_markup=ReplyKeyboardMarkup([["â¬…ï¸ Orqaga"]], one_time_keyboard=True, resize_keyboard=True)
    )
    return ADMIN_ACTION

def create_qr_code(data, output_path):
    try:
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=4,
        )
        qr.add_data(data)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        img.save(output_path)
        logger.info(f"QR kod saqlandi: {output_path}")
    except Exception as e:
        logger.error(f"QR kod yaratishda xato: {str(e)}")
        raise

async def add_to_queue(user_id: str, update: Update, context: ContextTypes.DEFAULT_TYPE) -> int:
    """Foydalanuvchi soâ€˜rovini navbatga qoâ€˜shadi"""
    try:
        await document_queue.put((user_id, update, context))
        queue_size = document_queue.qsize()
        logger.info(f"Soâ€˜rov navbatga qoâ€˜shildi: user_id={user_id}, navbat hajmi={queue_size}")
        await update.message.reply_text(f"â³ Soâ€˜rovingiz navbatda: {queue_size}-oâ€˜rinda. Iltimos, kuting...")
        return queue_size
    except Exception as e:
        logger.error(f"Navbatga qoâ€˜shishda xato: user_id={user_id}, {str(e)}")
        await update.message.reply_text("âŒ Soâ€˜rovni navbatga qoâ€˜shishda xato yuz berdi.")
        return -1

async def queue_worker():
    """Navbatdagi soâ€˜rovlarni ketma-ket bajaradi"""
    global queue_worker_running
    queue_worker_running = True
    logger.info("Navbat ishchisi ishga tushdi")
    while queue_worker_running:
        try:
            user_id, update, context = await document_queue.get()
            logger.info(f"Soâ€˜rov olindi: user_id={user_id}")
            try:
                async with context.bot_data.setdefault('queue_lock', asyncio.Lock()):
                    await generate_document_internal(update, context)
            except Exception as e:
                logger.error(f"Soâ€˜rovni bajarishda xato: user_id={user_id}, {str(e)}")
                await update.message.reply_text(f"âŒ Hujjatni yaratishda xato: {str(e)}")
            finally:
                document_queue.task_done()
                logger.info(f"Soâ€˜rov tugallandi: user_id={user_id}")
        except asyncio.CancelledError:
            logger.info("Navbat ishchisi toâ€˜xtatildi")
            break
        except Exception as e:
            logger.error(f"Navbat ishchisida xato: {str(e)}")
            await asyncio.sleep(1)
        
async def handle_qr_code(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)

    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['qr_code_data'] = None  
        context.user_data['qr_code_replacement'] = ""  
        await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
        return await generate_document(update, context)

    qr_data = update.message.text.strip()
    if qr_data == "ðŸ…° Avtomatik yasash":
        doc_type = context.user_data.get('doc_type')
        if doc_type in ['Sertifikat', 'Diplom']:
            qr_data = (
                f"{doc_type}\n"
                f"Taqdirlandi:{context.user_data.get('taqdirlangan', '')}\n"
                f"Taqdirladi:{context.user_data.get('taqdirlovchi', '')}\n"
                F"Sana:{context.user_data.get('sana', '')}"
            )
        elif doc_type == 'Shablon':
            qr_data = (
                f"{context.user_data.get('shablon_nomi', '')}: \n"
                f"Taqdirladi:{context.user_data.get('taqdirlovchi', '')}\n"
                F"Sana:{context.user_data.get('sana', '')}\n"
                f"Matn:{context.user_data.get('shablon_matni', '')}"
            )
    context.user_data['qr_code_data'] = qr_data
    context.user_data['qr_code_replacement'] = ""
    await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
    return await generate_document(update, context)

async def handle_template(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        selected_number = update.message.text.strip()
        if not selected_number.isdigit():
            await update.message.reply_text("Iltimos, faqat raqam kiriting.")
            return TEMPLATE
        selected_template = f"{selected_number}.pptx"
        doc_type = context.user_data.get('doc_type')
        if doc_type == 'Sertifikat':
            templates = context.user_data.get('templates', get_templates(SERTIFIKAT_PATH))
            if selected_template not in templates:
                await update.message.reply_text(f"âŒ {selected_number} raqamli namuna topilmadi. Iltimos, qaytadan tanlang.")
                return TEMPLATE
            template_path = os.path.join(SERTIFIKAT_PATH, selected_template)
            tags = SERTIFIKAT_TAGS
        elif doc_type == 'Taklifnoma':
            if int(selected_number) < 1 or int(selected_number) > 10:
                await update.message.reply_text("Iltimos, 1 dan 10 gacha raqam kiriting.")
                return TEMPLATE
            template_path = os.path.join(TAKLIFNOMA_PATH, selected_template)
            tags = TAKLIFNOMA_TAGS
        elif doc_type == 'Shablon':
            if int(selected_number) < 1 or int(selected_number) > 15:
                await update.message.reply_text("Iltimos, 1 dan 15 gacha raqam kiriting.")
                return TEMPLATE
            template_path = os.path.join(SHABLON_PATH, selected_template)
            tags = SHABLON_TAGS
        elif doc_type == 'Diplom':
            templates = context.user_data.get('templates', get_templates(DIPLOM_PATH))
            if int(selected_number) < 1 or int(selected_number) > 26:
                await update.message.reply_text("Iltimos, 1 dan 26 gacha raqam kiriting.")
                return TEMPLATE
            if selected_template not in templates:
                await update.message.reply_text(f"âŒ {selected_number} raqamli namuna topilmadi. Iltimos, qaytadan tanlang.")
                return TEMPLATE
            template_path = os.path.join(DIPLOM_PATH, selected_template)
            tags = DIPLOM_TAGS
        context.user_data['selected_template'] = template_path
        is_ready, msg, has_date_tag, present_tags, has_qr_code_tag = check_template(template_path, tags)
        if not is_ready:
            await update.message.reply_text(msg)
            return ConversationHandler.END
        context.user_data['has_date_tag'] = has_date_tag
        context.user_data['present_tags'] = present_tags
        context.user_data['has_qr_code_tag'] = has_qr_code_tag
        save_context_data(user_id, context.user_data)
        if doc_type in ['Sertifikat', 'Diplom']:
            if "{taqdirlangan}" in present_tags:
                user_data = load_user_data(user_id)
                last_taqdirlanuvchi = user_data.get('taqdirlangan', 'Bahodirov Dilxush')
                keyboard = [[last_taqdirlanuvchi]]
                await update.message.reply_text(
                    f"1ï¸âƒ£ Taqdirlanuvchi ismini kiriting.\n\n"
                    f"<b>Masalan:</b> <code>{last_taqdirlanuvchi}</code>\n"
                    "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz ismni kiriting:",
                    parse_mode="HTML",
                    reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return TAQDIRLANUVCHI
            else:
                context.user_data['taqdirlangan'] = ""
                if "{taqdirlovchi}" in present_tags:
                    user_data = load_user_data(user_id)
                    last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
                    keyboard = [[last_taqdirlovchi], ["â­ Oâ€˜tkazib yuborish"]]
                    await update.message.reply_text(
                        f"1ï¸âƒ£ Kim taqdirlayapti?\n\n"
                        f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code> \n"
                        "Iltimos, taqdirlovchi ismini yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                        parse_mode="HTML",
                        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                    return TAQDIRLOVCHI
                else:
                    context.user_data['taqdirlovchi'] = ""
                    if doc_type == 'Sertifikat' and "{sertifikat_matni}" in present_tags:
                        user_data = load_user_data(user_id)
                        last_sertifikat_matni = user_data.get('sertifikat_matni', 'Ushbu sertifikat bilan taqdirlanadi!')
                        keyboard = [["â­ Oâ€˜tkazib yuborish"]]
                        with open(INFOTEXT_IMAGE, "rb") as img5:
                            await update.message.reply_photo(
                                photo=img5,
                                caption=f"1ï¸âƒ£ Sertifikat matnini yuboring.\n\n"
                                        f"<i>Masalan:</i> \n"
                                        f"<code>{last_sertifikat_matni}</code> \n\n"
                                        "Yoki oâ€˜zingizga mos matn yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                        return SERTIFIKAT_MATNI
                    elif doc_type == 'Diplom' and "{diplom_matni}" in present_tags:
                        user_data = load_user_data(user_id)
                        last_diplom_matni = user_data.get('diplom_matni', 'Siz bizning kursimizni muvaffaqiyatli tamomlaganingiz uchun ushbu diplom bilan taqdirlanasiz!')
                        keyboard = [["â­ Oâ€˜tkazib yuborish"]]
                        with open(DIPLOM_MATNI_IMAGE, "rb") as img:
                            await update.message.reply_photo(
                                photo=img,
                                caption=f"1ï¸âƒ£ Diplom matnini yuboring.\n\n"
                                        f"<i>Masalan:</i> \n"
                                        f"<code>{last_diplom_matni}</code> \n\n"
                                        "Yoki oâ€˜zingizga mos matn yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                        return DIPLOM_MATNI
                    else:
                        context.user_data['sertifikat_matni'] = ""
                        context.user_data['diplom_matni'] = ""
                        if has_date_tag:
                            user_data = load_user_data(user_id)
                            last_date = user_data.get('sana', '14.10.2025')
                            keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                            await update.message.reply_text(
                                f"1ï¸âƒ£ {doc_type.capitalize()} uchun sana kiritingðŸ“†.\n\n"
                                f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                                "ðŸ“† Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                            return SANA
                        elif has_qr_code_tag:
                            keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
                            await update.message.reply_text(
                                f"ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                                f"<b>Masalan:</b> <code>https://example.com</code>\n"
                                "Yoki quyidagi tugmalardan birini tanlang:",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                            return QR_CODE
                        else:
                            context.user_data['sana'] = ""
                            await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                            return await generate_document(update, context)
        elif doc_type == 'Taklifnoma':
            if "{info_text}" in present_tags:
                user_data = load_user_data(user_id)
                last_info_text = user_data.get('info_text', (
                    "Hurmatli AZIZ MEHMONIMIZ\n Siz va oila a'zolaringizni 2025-yil 16-fevral kuni soat 18:00 da aziz farzandimiz\n FARIDJON va SHUKRONAXON larning\n Nikoh to'yi munosabai bilan yoziladigan dasturxonimizga lutgan taklif etamiz!"
                ))
                keyboard = [[last_info_text], ["â­ Oâ€˜tkazib yuborish"]]
                with open(TAKLIFNOMA_IMAGE, "rb") as img:
                    await update.message.reply_photo(
                        photo=img,
                        caption=f"1ï¸âƒ£ Taklifnoma matnini kiriting.\n\n"
                                f"<b>Masalan:</b> \n<code>{last_info_text}</code>\n"
                                "Matnni nusxalashingiz yoki oâ€˜zingiz yangi matn kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                        parse_mode="HTML",
                        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                return INFO_TEXT
            else:
                context.user_data['info_text'] = ""
                if "{manzil_va_ega}" in present_tags:
                    user_data = load_user_data(user_id)
                    last_manzil_va_ega = user_data.get('manzil_va_ega', 'Hurmat bilan Bahodirovlar oilasi\nManzil: Dehqonobod tumani Mamat ota toâ€˜yxonasi')
                    keyboard = [[last_manzil_va_ega], ["â­ Oâ€˜tkazib yuborish"]]
                    await update.message.reply_text(
                        f"1ï¸âƒ£ Manzil va ega ma'lumotlarini kiritingðŸ“Œ.\n\n"
                        f"<b>Masalan:</b> \n<code>{last_manzil_va_ega}</code>\n"
                        "Matnni nusxalashingiz yoki oâ€˜zingiz yangi matn kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                        parse_mode="HTML",
                        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                    return MANZIL_VA_EGA
                else:
                    context.user_data['manzil_va_ega'] = ""
                    if has_qr_code_tag:
                        keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
                        await update.message.reply_text(
                            f"ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                            f"<b>Masalan:</b> <code>https://example.com</code>\n"
                            "Yoki quyidagi tugmalardan birini tanlang:",
                            parse_mode="HTML",
                            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                        )
                        return QR_CODE
                    await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                    return await generate_document(update, context)
        elif doc_type == 'Shablon':
            if "{shablon_nomi}" in present_tags:
                user_data = load_user_data(user_id)
                last_shablon_nomi = user_data.get('shablon_nomi', 'TASHAKKURNOMA')
                keyboard = [[last_shablon_nomi]]
                with open(SHABLON_NOMI_IMAGE, "rb") as img:
                    await update.message.reply_photo(
                        photo=img,
                        caption=f"1ï¸âƒ£ Shablon sarlavhasini kiritingðŸ“‹.\n\n"
                                f"<b>Masalan:</b> <code>{last_shablon_nomi}</code>\n"
                                "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz nom kiriting:",
                        parse_mode="HTML",
                        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                return SHABLON_NOMI
            else:
                context.user_data['shablon_nomi'] = ""
                if "{shablon_matni}" in present_tags:
                    user_data = load_user_data(user_id)
                    last_shablon_matni = user_data.get('shablon_matni', 'Qashqadaryo viloiyati Dehqonobod tumani 87â€“sonli umumta\'lim maktabining 5-\'A\' sinf oâ€˜quvchisi Alisherov Farid 2024 â€“ 2025 O\'quv yilida a\'lo va yaxshi baholari hamda namunali xulqi uchun MAQTOV YORLIG\'I bilan taqdirlanadi')
                    keyboard = [[last_shablon_matni], ["â­ Oâ€˜tkazib yuborish"]]
                    with open(SHABLON_MATNI_IMAGE, "rb") as img:
                        await update.message.reply_photo(
                            photo=img,
                            caption=f"1ï¸âƒ£ Shablon matnini kiritingðŸ“‹.\n\n"
                                    f"<b>Masalan:</b> \n<code>{last_shablon_matni}</code>\n"
                                    "Matnni nusxalashingiz yoki oâ€˜zingiz yangi matn kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                            parse_mode="HTML",
                            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                        )
                    return SHABLON_MATNI
                else:
                    context.user_data['shablon_matni'] = ""
                    if "{taqdirlovchi}" in present_tags:
                        user_data = load_user_data(user_id)
                        last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
                        keyboard = [[last_taqdirlovchi], ["â­ Oâ€˜tkazib yuborish"]]
                        await update.message.reply_text(
                            f"1ï¸âƒ£ Taqdirlovchi ismini kiritingðŸ‘¤.\n\n"
                            f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code>\n"
                            "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz ismini kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                            parse_mode="HTML",
                            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                        )
                        return TAQDIRLOVCHI
                    else:
                        context.user_data['taqdirlovchi'] = ""
                        if has_date_tag:
                            user_data = load_user_data(user_id)
                            last_date = user_data.get('sana', '14.10.2025')
                            keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                            await update.message.reply_text(
                                f"1ï¸âƒ£ Shablon uchun sana kiritingðŸ“†.\n\n"
                                f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                                "ðŸ“† Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                            return SANA
                        elif has_qr_code_tag:
                            keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
                            await update.message.reply_text(
                                f"ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                                f"<b>Masalan:</b> <code>https://example.com</code>\n"
                                "Yoki quyidagi tugmalardan birini tanlang:",
                                parse_mode="HTML",
                                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                            )
                            return QR_CODE
                        else:
                            context.user_data['sana'] = ""
                            await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                            return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return TEMPLATE
    except Exception as e:
        logger.error(f"Shablonni tanlashda xato: {str(e)}")
        await update.message.reply_text("âŒ Shablonni tanlashda xatolik yuz berdi. Iltimos, qaytadan urinib koâ€˜ring.")
        return TEMPLATE

async def handle_diplom_matni(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['diplom_matni'] = ""
    else:
        diplom_matni = update.message.text.strip()
        user_data = load_user_data(user_id)
        user_data['diplom_matni'] = diplom_matni
        save_user_data(user_id, user_data)
        context.user_data['diplom_matni'] = diplom_matni
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="Diplom matni va taqdirlovchi uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return SHRIFT2

async def handle_shablon(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        template_num = int(update.message.text.strip())
        templates = get_templates(SHABLON_PATH)
        if template_num < 1 or template_num > len(templates):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shablon raqami. Iltimos, qaytadan tanlang.")
            return SHABLON
        template_path = os.path.join(SHABLON_PATH, templates[template_num - 1])
        is_valid, message, has_date_tag, present_tags, has_qr_code_tag = check_template(template_path, SHABLON_TAGS)
        if not is_valid:
            await update.message.reply_text(message)
            return SHABLON
        context.user_data['selected_template'] = template_path
        context.user_data['has_date_tag'] = has_date_tag
        context.user_data['present_tags'] = present_tags
        context.user_data['has_qr_code_tag'] = has_qr_code_tag
        user_data = load_user_data(user_id)
        last_shablon_nomi = user_data.get('shablon_nomi', 'TASHAKKURNOMA')
        keyboard = [[last_shablon_nomi], ["â­ Oâ€˜tkazib yuborish"]]
        with open(SHABLON_NOMI_IMAGE, "rb") as img:
            await update.message.reply_photo(
                photo=img,
                caption=f"1ï¸âƒ£ Shablon nomini kiritingðŸ“.\n\n"
                        f"<b>Masalan:</b> <code>{last_shablon_nomi}</code>\n"
                        "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz nom kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        return SHABLON_NOMI
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHABLON

async def handle_shablon_nomi(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    shablon_nomi = update.message.text.strip()
    user_data = load_user_data(user_id)
    user_data['shablon_nomi'] = shablon_nomi
    user_data['username'] = update.effective_user.username or "Nomaâ€˜lum"
    save_user_data(user_id, user_data)
    context.user_data['shablon_nomi'] = shablon_nomi
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="Shablon sarlavhasi uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return SHABLON_SHRIFT1

async def handle_shablon_shrift1(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return SHABLON_SHRIFT1
        context.user_data['shablon_nomi_font'] = FONTS[font_index]
        if "{shablon_matni}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_shablon_matni = user_data.get('shablon_matni', 'Qashqadaryo viloiyati Dehqonobod tumani 87â€“sonli umumtalim maktabining 5-A sinf oâ€˜quvchisi Alisherov Farid 2024 â€“ 2025 uquv yilida alo va yaxshi baholari hamda namunali xulqi uchun MAQTOV YORLIGI bilan taqdirlanadi')
            keyboard = [[last_shablon_matni], ["â­ Oâ€˜tkazib yuborish"]]
            with open(SHABLON_MATNI_IMAGE, "rb") as img:
                await update.message.reply_photo(
                    photo=img,
                    caption=f"2ï¸âƒ£ Shablon matnini kiritingðŸ“‹.\n\n"
                            f"<b>Masalan:</b> \n<code>{last_shablon_matni}</code>\n"
                            "Matnni nusxalashingiz yoki oâ€˜zingiz yangi matn kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
            return SHABLON_MATNI
        else:
            context.user_data['shablon_matni'] = ""
            if "{taqdirlovchi}" in context.user_data.get('present_tags', []):
                user_data = load_user_data(user_id)
                last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
                keyboard = [[last_taqdirlovchi]]
                await update.message.reply_text(
                    f"2ï¸âƒ£ Taqdirlovchi ismini kiritingðŸ‘¤.\n\n"
                    f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code>\n"
                    "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz ismini kiriting:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return TAQDIRLOVCHI
            else:
                context.user_data['taqdirlovchi'] = ""
                if context.user_data.get('has_date_tag'):
                    user_data = load_user_data(user_id)
                    last_date = user_data.get('sana', '14.10.2025')
                    keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                    await update.message.reply_text(
                        f"2ï¸âƒ£ Shablon uchun sana kiritingðŸ“†.\n\n"
                        f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                        "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                        parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                    return SANA
                else:
                    context.user_data['sana'] = ""
                    await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                    return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHABLON_SHRIFT1

async def handle_shablon_matni(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['shablon_matni'] = ""
        context.user_data['taqdirlovchi'] = ""
        context.user_data['sana'] = ""
        if "{taqdirlovchi}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
            keyboard = [[last_taqdirlovchi]]
            await update.message.reply_text(
                f"2ï¸âƒ£ Taqdirlovchi ismini kiritingðŸ‘¤.\n\n"
                f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code>\n"
                "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz ismini kiriting:",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return TAQDIRLOVCHI
        else:
            context.user_data['taqdirlovchi'] = ""
            if context.user_data.get('has_date_tag'):
                user_data = load_user_data(user_id)
                last_date = user_data.get('sana', '14.10.2025')
                keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                await update.message.reply_text(
                    f"2ï¸âƒ£ Shablon uchun sana kiritingðŸ“†.\n\n"
                    f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                    "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return SANA
            else:
                context.user_data['sana'] = ""
                await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                return await generate_document(update, context)
    shablon_matni = update.message.text.strip()
    user_data = load_user_data(user_id)
    user_data['shablon_matni'] = shablon_matni
    save_user_data(user_id, user_data)
    context.user_data['shablon_matni'] = shablon_matni
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="Shablon matni uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return SHABLON_SHRIFT2

async def handle_shablon_shrift2(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return SHABLON_SHRIFT2

        context.user_data['shablon_matni_font'] = FONTS[font_index]
        context.user_data['shablon_sana_font'] = FONTS[font_index]

        if "{taqdirlovchi}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
            keyboard = [[last_taqdirlovchi], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                f"3ï¸âƒ£ Taqdirlovchi ismini kiritingðŸ‘¤.\n\n"
                f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code>\n"
                "Yuqoridagi matnni tanlab nusxalashingiz yoki oâ€˜zingiz ismini kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return TAQDIRLOVCHI
        else:
            context.user_data['taqdirlovchi'] = ""
            if context.user_data.get('has_date_tag'):
                user_data = load_user_data(user_id)
                last_date = user_data.get('sana', '14.10.2025')
                keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                await update.message.reply_text(
                    f"3ï¸âƒ£ Shablon uchun sana kiritingðŸ“†.\n\n"
                    f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                    "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return SANA
            else:
                context.user_data['sana'] = ""
                await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHABLON_SHRIFT2

async def handle_taqdirlanuvchi(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    taqdirlanuvchi = update.message.text.strip()
    user_data = load_user_data(user_id)
    user_data['taqdirlangan'] = taqdirlanuvchi
    user_data['username'] = update.effective_user.username or "Nomaâ€˜lum"
    save_user_data(user_id, user_data)
    context.user_data['taqdirlangan'] = taqdirlanuvchi
    save_context_data(user_id, context.user_data)  
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="ðŸ‘¤Taqdirlanuvchi ismi uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return SHRIFT1

async def handle_shrift1(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return SHRIFT1
        context.user_data['taqdirlanuvchi_font'] = FONTS[font_index]
        if "{taqdirlovchi}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_taqdirlovchi = user_data.get('taqdirlovchi', 'Bahodirov Dilxush')
            keyboard = [[last_taqdirlovchi], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                f"2ï¸âƒ£ Kim taqdirlayapti?ðŸ‘¤\n\n"
                f"<b>Masalan:</b> <code>{last_taqdirlovchi}</code> \n"
                "Iltimos,ðŸ‘¤Taqdirlovchi ismini yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return TAQDIRLOVCHI
        else:
            context.user_data['taqdirlovchi'] = ""
            if "{sertifikat_matni}" in context.user_data.get('present_tags', []):
                user_data = load_user_data(user_id)
                last_sertifikat_matni = user_data.get('sertifikat_matni')
                keyboard = [["â­ Oâ€˜tkazib yuborish"]]
                with open(INFOTEXT_IMAGE, "rb") as img5:
                    await update.message.reply_photo(
                        photo=img5,
                        caption=f"2ï¸âƒ£ Sertifikat matnini yuboringðŸ“.\n\n"
                                f"<i>Masalan:</i> \n"
                                f"<code>{last_sertifikat_matni}</code> \n\n"
                                "Yoki oâ€˜zingizga mos matn yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                        parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                return SERTIFIKAT_MATNI
            else:
                context.user_data['sertifikat_matni'] = ""
                if context.user_data.get('has_date_tag'):
                    user_data = load_user_data(user_id)
                    last_date = user_data.get('sana', '14.10.2025')
                    keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                    await update.message.reply_text(
                        f"2ï¸âƒ£ Sertifikat uchun sana kiritingðŸ“†.\n\n"
                        f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                        "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                        parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                    )
                    return SANA
                else:
                    context.user_data['sana'] = ""
                    await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                    return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHRIFT1

async def handle_taqdirlovchi(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['taqdirlovchi'] = ""
    else:
        taqdirlovchi = update.message.text.strip()
        user_data = load_user_data(user_id)
        user_data['taqdirlovchi'] = taqdirlovchi
        save_user_data(user_id, user_data)
        context.user_data['taqdirlovchi'] = taqdirlovchi
    if context.user_data.get('doc_type') == 'Sertifikat':
        if "{sertifikat_matni}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_sertifikat_matni = user_data.get('sertifikat_matni')
            keyboard = [["â­ Oâ€˜tkazib yuborish"]]
            with open(INFOTEXT_IMAGE, "rb") as img5:
                await update.message.reply_photo(
                    photo=img5,
                    caption=f"3ï¸âƒ£ Sertifikat matnini yuboringðŸ“.\n\n"
                            f"<i>Masalan:</i> \n"
                            f"<code>{last_sertifikat_matni}</code> \n\n"
                            "Yoki oâ€˜zingizga mos matn yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
            return SERTIFIKAT_MATNI
        else:
            context.user_data['sertifikat_matni'] = ""
            if context.user_data.get('has_date_tag'):
                user_data = load_user_data(user_id)
                last_date = user_data.get('sana', '14.10.2025')
                keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                await update.message.reply_text(
                    f"3ï¸âƒ£ Sertifikat uchun sana kiritingðŸ“†.\n\n"
                    f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                    "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return SANA
            else:
                context.user_data['sana'] = ""
                await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                return await generate_document(update, context)
    elif context.user_data.get('doc_type') == 'Diplom':
        if "{diplom_matni}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_diplom_matni = user_data.get('diplom_matni', 'Siz bizning kursimizni muvaffaqiyatli tamomlaganingiz uchun ushbu diplom bilan taqdirlanasiz!')
            keyboard = [["â­ Oâ€˜tkazib yuborish"]]
            with open(DIPLOM_MATNI_IMAGE, "rb") as img:
                await update.message.reply_photo(
                    photo=img,
                    caption=f"3ï¸âƒ£ Diplom matnini yuboringðŸ“.\n\n"
                            f"<i>Masalan:</i> \n"
                            f"<code>{last_diplom_matni}</code> \n\n"
                            "Yoki oâ€˜zingizga mos matn yuboring yoki oâ€˜tkazib yuborishni tanlang.",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
            return DIPLOM_MATNI
        else:
            context.user_data['diplom_matni'] = ""
            if context.user_data.get('has_date_tag'):
                user_data = load_user_data(user_id)
                last_date = user_data.get('sana', '14.10.2025')
                keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
                await update.message.reply_text(
                    f"3ï¸âƒ£ Diplom uchun sana kiritingðŸ“†.\n\n"
                    f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                    "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                    parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
                )
                return SANA
            else:
                context.user_data['sana'] = ""
                await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
                return await generate_document(update, context)
    else:
        if context.user_data.get('has_date_tag'):
            user_data = load_user_data(user_id)
            last_date = user_data.get('sana', '14.10.2025')
            keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                f"4ï¸âƒ£ Shablon uchun sana kiritingðŸ“†.\n\n"
                f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return SANA
        else:
            context.user_data['sana'] = ""
            await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
            return await generate_document(update, context)

async def handle_sertifikat_matni(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['sertifikat_matni'] = ""
    else:
        sertifikat_matni = update.message.text.strip()
        user_data = load_user_data(user_id)
        user_data['sertifikat_matni'] = sertifikat_matni
        save_user_data(user_id, user_data)
        context.user_data['sertifikat_matni'] = sertifikat_matni
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="Sertifikat matni va taqdirlovchi uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return SHRIFT2

async def handle_shrift2(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return SHRIFT2
        context.user_data['other_font'] = FONTS[font_index]
        if context.user_data.get('has_date_tag'):
            user_data = load_user_data(user_id)
            last_date = user_data.get('sana', '14.10.2025')
            keyboard = [["Bugungi sana"], [last_date], ["â­ Oâ€˜tkazib yuborish"]] if user_data.get('sana') else [["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                f"4ï¸âƒ£ Sertifikat uchun sana kiritingðŸ“†.\n\n"
                f"<b>Masalan:</b> <code>14.10.2025</code>\n"
                "ðŸ“†Sanani DD.MM.YYYY formatida kiriting yoki tugmalardan birini tanlang:",
                parse_mode="HTML", reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return SANA
        elif context.user_data.get('has_qr_code_tag'):
            keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                "ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                "<b>Masalan:</b> <code>https://example.com</code>\n"
                "Yoki quyidagi tugmalardan birini tanlang:",
                parse_mode="HTML",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return QR_CODE
        else:
            context.user_data['sana'] = ""
            await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
            return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHRIFT2

async def handle_sana(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    if update.message.text == "â­ Oâ€˜tkazib yuborish":
        context.user_data['sana'] = ""
    elif update.message.text == "Bugungi sana":
        context.user_data['sana'] = datetime.datetime.now().strftime("%d.%m.%Y")
    else:
        date_text = update.message.text.strip()
        if not validate_date(date_text):
            await update.message.reply_text(
                "âŒ Notoâ€˜gâ€˜ri sana formati.ðŸ“† Iltimos, DD.MM.YYYY formatida kiriting (masalan, 14.10.2025).",
                reply_markup=ReplyKeyboardMarkup([["Bugungi sana"], ["â­ Oâ€˜tkazib yuborish"]], one_time_keyboard=True, resize_keyboard=True)
            )
            return SANA
        context.user_data['sana'] = date_text
    user_data = load_user_data(user_id)
    user_data['sana'] = context.user_data['sana']
    save_user_data(user_id, user_data)
    if context.user_data.get('doc_type') == 'shablon' and context.user_data.get('has_date_tag'):
        if not os.path.exists(SHRIFT_IMAGE):
            await update.message.reply_text("âŒ shrift.jpg fayli topilmadi.")
            return ConversationHandler.END
        keyboard = create_font_keyboard()
        with open(SHRIFT_IMAGE, "rb") as img:
            await update.message.reply_photo(
                photo=img,
                caption="ðŸ‘¤Taqdirlovchi va Sana uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
        return SHABLON_SHRIFT3
    elif context.user_data.get('has_qr_code_tag'):
        keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
        await update.message.reply_text(
            "ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
            "<b>Masalan:</b> <code>https://example.com</code>\n"
            "Yoki quyidagi tugmalardan birini tanlang:",
            parse_mode="HTML",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
        return QR_CODE
    await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
    return await generate_document(update, context)

async def handle_shablon_shrift3(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return SHABLON_SHRIFT3
        context.user_data['shablon_sana_font'] = FONTS[font_index]
        if context.user_data.get('has_qr_code_tag'):
            keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                "ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                "<b>Masalan:</b> <code>https://example.com</code>\n"
                "Yoki quyidagi tugmalardan birini tanlang:",
                parse_mode="HTML",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return QR_CODE
        await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
        return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return SHABLON_SHRIFT3

async def handle_info_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)

    info_text = update.message.text.strip()
    user_data = load_user_data(user_id)
    user_data['info_text'] = info_text
    save_user_data(user_id, user_data)
    context.user_data['info_text'] = info_text

    if not os.path.exists(SHRIFT_IMAGE):
        await update.message.reply_text("âŒ shrift.jpg fayli topilmadi.")
        return ConversationHandler.END

    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="ðŸ’ŒTaklifnoma matni uchun shriftni tanlang:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return TAKLIFNOMA_SHRIFT1

async def handle_taklifnoma_shrift1(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)

    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return TAKLIFNOMA_SHRIFT1

        context.user_data['info_text_font'] = FONTS[font_index]
        logger.info(f"info_text_font saqlandi: {FONTS[font_index]}")  # Log qoâ€˜shildi

        if "{manzil_va_ega}" in context.user_data.get('present_tags', []):
            user_data = load_user_data(user_id)
            last_manzil_va_ega = user_data.get('manzil_va_ega', 'Hurmat bilan Bahodirovlar oilasi\nManzil: Dehqonobod tumani Mamat ota toâ€˜yxonasi')
            keyboard = [[last_manzil_va_ega], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                f"2ï¸âƒ£ Manzil va ega ma'lumotlarini kiritingðŸ“Œ.\n\n"
                f"<b>Masalan:</b> \n<code>{last_manzil_va_ega}</code>\n"
                "Matnni nusxalashingiz yoki oâ€˜zingiz yangi matn kiriting yoki oâ€˜tkazib yuborishni tanlang:",
                parse_mode="HTML",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return MANZIL_VA_EGA
        else:
            context.user_data['manzil_va_ega'] = ""
            await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
            return await generate_document(update, context)

    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return TAKLIFNOMA_SHRIFT1

async def handle_manzil_va_ega(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    manzil_va_ega = update.message.text.strip()
    user_data = load_user_data(user_id)
    user_data['manzil_va_ega'] = manzil_va_ega
    save_user_data(user_id, user_data)
    context.user_data['manzil_va_ega'] = manzil_va_ega
    keyboard = create_font_keyboard()
    with open(SHRIFT_IMAGE, "rb") as img:
        await update.message.reply_photo(
            photo=img,
            caption="Manzil va ega uchun shriftni tanlangðŸ“Œ:\nQuyidagi raqamlardan birini tanlang:",
            reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
        )
    return TAKLIFNOMA_SHRIFT2

async def handle_taklifnoma_shrift2(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    try:
        font_index = int(update.message.text.strip()) - 1
        if font_index < 0 or font_index >= len(FONTS):
            await update.message.reply_text("âŒ Notoâ€˜gâ€˜ri shrift raqami. Iltimos, qaytadan tanlang.")
            return TAKLIFNOMA_SHRIFT2
        context.user_data['manzil_va_ega_font'] = FONTS[font_index]
        logger.info(f"manzil_va_ega_font saqlandi: {FONTS[font_index]}")  # Log qoâ€˜shildi
        if context.user_data.get('has_qr_code_tag'):
            keyboard = [["ðŸ…° Avtomatik yasash"], ["â­ Oâ€˜tkazib yuborish"]]
            await update.message.reply_text(
                "ðŸ“· QR kod mazmuni ma'lumotini kiriting:\n\n"
                "<b>Masalan:</b> <code>https://example.com</code>\n"
                "Yoki quyidagi tugmalardan birini tanlang:",
                parse_mode="HTML",
                reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
            )
            return QR_CODE
        await update.message.reply_text("â³", reply_markup=ReplyKeyboardRemove())
        return await generate_document(update, context)
    except ValueError:
        await update.message.reply_text("Iltimos, faqat raqam kiriting.")
        return TAKLIFNOMA_SHRIFT2


async def generate_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Hujjat yaratish so'rovini navbatga qo'shadi
    """
    user_id = update.effective_user.id
    return await add_to_queue(user_id, update, context)

async def generate_document_internal(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    Foydalanuvchi uchun hujjat (masalan, Taklifnoma) yaratadi va JPG sifatida yuboradi.
    """
    user_id = update.effective_user.id
    config = load_config()  # Konfiguratsiya faylini yuklash
    user_data = load_user_data(user_id)  # Foydalanuvchi ma'lumotlarini yuklash
    doc_type = context.user_data.get('doc_type', '')
    template_path = context.user_data.get('selected_template')
    logger.info(f"Generating document for user {user_id}, doc_type: {doc_type}, template: {template_path}")
    logger.debug(f"context.user_data: {context.user_data}")

    # Hujjat narxini aniqlash
    cost = {
        'Sertifikat': config['CERTIFICATE_COST'],
        'Taklifnoma': config['TAKLIFNOMA_COST'],
        'Shablon': config['SHABLON_COST'],
        'Diplom': config['DIPLOM_COST']
    }.get(doc_type, 0)

    # Balansni tekshirish
    if user_data.get('balance', 5000) < cost:
        await update.message.reply_text(
            f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q.ðŸ’µJoriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
            f"{doc_type.capitalize()} narxi: {cost} soâ€˜m.\nIltimos, balansingizni toâ€˜ldiring.",
            reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return BALANCE

    # Temp papkasini yaratish
    if not os.path.exists(TEMP_PATH):
        os.makedirs(TEMP_PATH)
        logger.info(f"Created temp directory: {TEMP_PATH}")

    async def process_presentation():
        logger.info(f"Starting presentation processing for user {user_id}")
        try:
            prs = Presentation(template_path)
            logger.info(f"Loaded presentation from {template_path}")
            replacements = {
                "{qr_code}": context.user_data.get('qr_code_replacement', "")
            }
            qr_code_path = None
            if context.user_data.get('has_qr_code_tag') and context.user_data.get('qr_code_data'): 
                qr_code_path = os.path.join(TEMP_PATH, f"qr_code_{user_id}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S_%f')}.png")
                create_qr_code(context.user_data['qr_code_data'], qr_code_path)
                logger.info(f"Generated QR code for user {user_id} at {qr_code_path}")

            if doc_type == 'Taklifnoma':
                replacements.update({
                    "{info_text}": context.user_data.get('info_text', ""),
                    "{manzil_va_ega}": context.user_data.get('manzil_va_ega', "")
                })
                replace_text_and_font(
                    prs,
                    replacements,
                    info_text_font=context.user_data.get('info_text_font'),
                    manzil_va_ega_font=context.user_data.get('manzil_va_ega_font'),
                    qr_code_path=qr_code_path
                )

            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f")
            output_pptx = os.path.join(TEMP_PATH, f"output_{user_id}_{timestamp}.pptx")
            output_jpg = os.path.join(TEMP_PATH, f"output_{user_id}_{timestamp}.jpg")
            prs.save(output_pptx)
            logger.info(f"Saved PPTX to {output_pptx}")

            if not await convert_pptx_to_jpg(output_pptx, output_jpg, context):
                logger.error(f"Failed to convert PPTX to JPG: {output_jpg}")
                raise Exception("JPG faylini yaratishda xatolik yuz berdi")

            if qr_code_path and os.path.exists(qr_code_path):
                os.remove(qr_code_path)
                logger.info(f"Removed QR code file: {qr_code_path}")

            if not os.path.exists(output_jpg):
                logger.error(f"JPG fayli topilmadi: {output_jpg}")
                raise Exception("JPG fayli yaratilmadi")

            return output_pptx, output_jpg

        except Exception as e:
            logger.error(f"Presentation processing failed: {str(e)}", exc_info=True)
            raise

    try:
        output_pptx, output_jpg = await process_presentation()
        user_data['balance'] -= cost
        user_data['result_count'] = user_data.get('result_count', 0) + 1
        save_user_data(user_id, user_data)

        try:
            with open(output_jpg, "rb") as img:
                caption = "âœ… Taklifnomangiz tayyor boâ€˜ldi!"
                logger.info(f"Sending photo to user {user_id}: {output_jpg}")
                await update.message.reply_photo(
                    photo=img,
                    caption=caption,
                    reply_markup=ReplyKeyboardRemove()
                )
                logger.info(f"Photo sent successfully to user {user_id}")
        finally:
            with file_lock:  # Fayl blokirovkasi
                if os.path.exists(output_pptx):
                    os.remove(output_pptx)
                    logger.info(f"Removed PPTX file: {output_pptx}")
                if os.path.exists(output_jpg):
                    os.remove(output_jpg)
                    logger.info(f"Removed JPG file: {output_jpg}")

        await update.message.reply_text(
            f"ðŸ’¸Balansingizdan {cost} soâ€˜m yechildi.\n ðŸ’µJoriy balans: {user_data['balance']} soâ€˜m\n"
            f"ðŸ“ŠUmumiy natijalar soni: {user_data['result_count']}\n"
            "Natijani PDF formatda olishni xohlaysizmi?",
            reply_markup=ReplyKeyboardMarkup(
                [["Ha"], ["Yoâ€˜q"], ["Bosh menyuga qaytish"]], one_time_keyboard=True, resize_keyboard=True)
        )

        context.user_data['output_pptx'] = output_pptx
        context.user_data['output_jpg'] = output_jpg
        return PDF_CONFIRM

    except Exception as e:
        logger.error(f"Error generating document for user {user_id}: {str(e)}", exc_info=True)
        try:
            with file_lock:
                if 'output_pptx' in locals() and os.path.exists(output_pptx):
                    os.remove(output_pptx)
                    logger.info(f"Removed PPTX file: {output_pptx}")
                if 'output_jpg' in locals() and os.path.exists(output_jpg):
                    os.remove(output_jpg)
                    logger.info(f"Removed JPG file: {output_jpg}")
        except NameError:
            pass
        await update.message.reply_text(
            f"âŒ Hujjatni tayyorlashda xatolik yuz berdi: {str(e)}",
            reply_markup=ReplyKeyboardMarkup(
                [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸ”µ Taklifnoma yasash"],
                 ["ðŸŸ¢ Namuna", "ðŸŸ¤ Diplom yasash"],
                 ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                 ["ðŸ“ž Admin bilan aloqa"]], one_time_keyboard=True, resize_keyboard=True)
        )
        return START

async def handle_pdf_confirm(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = str(update.effective_user.id)
    config = load_config()
    doc_type = context.user_data.get('doc_type', '')
    user_data = load_user_data(user_id)
    text = update.message.text.strip() if update.message else ""

    # JPG faylini tozalash
    output_jpg_path = context.user_data.get('output_jpg', '')
    if os.path.exists(output_jpg_path):
        try:
            os.remove(output_jpg_path)
            logger.info(f"Removed JPG file: {output_jpg_path}")
        except Exception as e:
            logger.error(f"Error removing JPG file {output_jpg_path}: {str(e)}")

    # "Bosh menyuga qaytish" or "Yoâ€˜q" ni qayta ishlash
    if text in ["Bosh menyuga qaytish", "Yoâ€˜q"]:
        pdf_path = context.user_data.get('output_pdf', '')
        if pdf_path and os.path.exists(pdf_path):
            try:
                os.remove(pdf_path)
                logger.info(f"Removed PDF file: {pdf_path}")
                context.user_data.pop('output_pdf', None)
            except Exception as e:
                logger.error(f"Error removing PDF file {pdf_path}: {str(e)}")
        await update.message.reply_text(
            "âœ… PDF soâ€˜rovi bekor qilindi." if text == "Yoâ€˜q" else "Asosiy menyuga qaytdingiz.",
            reply_markup=ReplyKeyboardMarkup(
                [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                 ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                 ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                 ["ðŸ“ž Biz bilan aloqa"]],
                one_time_keyboard=True, resize_keyboard=True
            )
        )
        return START

    # "Ha" ni qayta ishlash (PDF yaratish)
    if text == "Ha":
        # Balansni tekshirish
        if user_data.get('balance', 5000) < config['PDF_COST']:
            await update.message.reply_text(
                f"âŒ Balansingizda yetarli mablagâ€˜ yoâ€˜q. ðŸ’µJoriy balans: {user_data.get('balance', 5000)} soâ€˜m. "
                f"PDF narxi: {config['PDF_COST']} soâ€˜m.\nIltimos, balansingizni toâ€˜ldiring.",
                reply_markup=ReplyKeyboardMarkup([["ðŸ’¸ Balansni toâ€˜ldirish"]], one_time_keyboard=True, resize_keyboard=True)
            )
            return BALANCE

        try:
            # Vaqt belgisi va fayl yoâ€˜llarini yaratish
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_pptx = os.path.join(TEMP_PATH, f"output_{user_id}_{timestamp}.pptx")
            output_pdf = os.path.join(TEMP_PATH, f"output_{user_id}_{timestamp}.pdf")
            logger.info(f"Generating PDF for user {user_id}: {output_pdf}")

            # PPTX faylini foydalanuvchi ma'lumotlari bilan qayta yaratish
            prs = Presentation(context.user_data.get('selected_template'))
            replacements = {
                "{qr_code}": context.user_data.get('qr_code_replacement', "")
            }
            qr_code_path = None
            if context.user_data.get('has_qr_code_tag') and context.user_data.get('qr_code_data'):
                qr_code_path = os.path.join(TEMP_PATH, f"qr_code_{user_id}_{timestamp}.png")
                create_qr_code(context.user_data['qr_code_data'], qr_code_path)
                logger.info(f"Generated QR code at {qr_code_path}")

            # Hujjat turiga qarab almashtirishlarni amalga oshirish
            if doc_type == 'Sertifikat':
                replacements.update({
                    "{taqdirlangan}": context.user_data.get('taqdirlangan', ""),
                    "{taqdirlovchi}": context.user_data.get('taqdirlovchi', ""),
                    "{sertifikat_matni}": context.user_data.get('sertifikat_matni', ""),
                    "{sana}": context.user_data.get('sana', "")
                })
                success = await replace_text_and_font(
                    prs,
                    replacements,
                    taqdirlanuvchi_font=context.user_data.get('taqdirlanuvchi_font'),
                    other_font=context.user_data.get('other_font'),
                    sana_font=context.user_data.get('other_font'),
                    qr_code_path=qr_code_path
                )
                if not success:
                    raise Exception("Failed to replace text and font for Sertifikat")
            elif doc_type == 'Taklifnoma':
                replacements.update({
                    "{info_text}": context.user_data.get('info_text', ""),
                    "{manzil_va_ega}": context.user_data.get('manzil_va_ega', "")
                })
                success = await replace_text_and_font(
                    prs,
                    replacements,
                    info_text_font=context.user_data.get('info_text_font'),
                    manzil_va_ega_font=context.user_data.get('manzil_va_ega_font'),
                    qr_code_path=qr_code_path
                )
                if not success:
                    raise Exception("Failed to replace text and font for Taklifnoma")
            elif doc_type == 'Shablon':
                replacements.update({
                    "{shablon_nomi}": context.user_data.get('shablon_nomi', ""),
                    "{shablon_matni}": context.user_data.get('shablon_matni', ""),
                    "{taqdirlovchi}": context.user_data.get('taqdirlovchi', ""),
                    "{sana}": context.user_data.get('sana', "")
                })
                success = await replace_text_and_font(
                    prs,
                    replacements,
                    shablon_nomi_font=context.user_data.get('shablon_nomi_font'),
                    shablon_matni_font=context.user_data.get('shablon_matni_font'),
                    shablon_sana_font=context.user_data.get('shablon_sana_font'),
                    qr_code_path=qr_code_path
                )
                if not success:
                    raise Exception("Failed to replace text and font for Shablon")
            elif doc_type == 'Diplom':
                replacements.update({
                    "{taqdirlangan}": context.user_data.get('taqdirlangan', ""),
                    "{taqdirlovchi}": context.user_data.get('taqdirlovchi', ""),
                    "{diplom_matni}": context.user_data.get('diplom_matni', ""),
                    "{sana}": context.user_data.get('sana', "")
                })
                success = await replace_text_and_font(
                    prs,
                    replacements,
                    taqdirlanuvchi_font=context.user_data.get('taqdirlanuvchi_font'),
                    other_font=context.user_data.get('other_font'),
                    sana_font=context.user_data.get('other_font'),
                    diplom_matni_font=context.user_data.get('other_font'),
                    qr_code_path=qr_code_path
                )
                if not success:
                    raise Exception("Failed to replace text and font for Diplom")

            # PPTX faylini saqlash
            prs.save(output_pptx)
            logger.info(f"Saved PPTX to {output_pptx}")

            # PPTX ni PDF ga aylantirish
            success = await convert_pptx_to_pdf(output_pptx, output_pdf, context)
            if not success:
                raise Exception("PDF faylini yaratishda xatolik yuz berdi")

            logger.info(f"Converted PPTX to PDF: {output_pdf}")

            # Vaqtinchalik fayllarni tozalash
            with file_lock:
                for file_path in [output_pptx, qr_code_path]:
                    if file_path and os.path.exists(file_path):
                        try:
                            os.remove(file_path)
                            logger.info(f"Removed file: {file_path}")
                        except Exception as e:
                            logger.error(f"Error removing file {file_path}: {str(e)}")

            # Balansdan PDF narxini yechish
            user_data['balance'] -= config['PDF_COST']
            user_data['result_count'] = user_data.get('result_count', 0) + 1
            save_user_data(user_id, user_data)
            logger.info(f"Deducted {config['PDF_COST']} from user {user_id}'s balance. New balance: {user_data['balance']}")

            # PDF fayl nomini hujjat turiga qarab belgilash
            pdf_filename = "Natija.pdf"
            if doc_type == 'Sertifikat':
                taqdirlanuvchi = context.user_data.get('taqdirlangan', '').replace(" ", "_")
                pdf_filename = f"Sertifikat_{taqdirlanuvchi}.pdf"
            elif doc_type == 'Taklifnoma':
                pdf_filename = "Taklifnoma.pdf"
            elif doc_type == 'Shablon':
                pdf_filename = "Shablon.pdf"
            elif doc_type == 'Diplom':
                taqdirlanuvchi = context.user_data.get('taqdirlangan', '').replace(" ", "_")
                pdf_filename = f"Diplom_{taqdirlanuvchi}.pdf"

            # Hujjat turiga qarab caption belgilash
            caption = "âœ… Natijangiz PDF formatda tayyor boâ€˜ldi!"
            if doc_type == 'Sertifikat':
                caption = "âœ… Sertifikatingiz PDF formatda tayyor boâ€˜ldi!"
            elif doc_type == 'Taklifnoma':
                caption = "âœ… Taklifnomangiz PDF formatda tayyor boâ€˜ldi!"
            elif doc_type == 'Diplom':
                caption = "âœ… Diplomingiz PDF formatda tayyor boâ€˜ldi!"

            # PDF faylini foydalanuvchiga yuborish
            try:
                with open(output_pdf, "rb") as pdf:
                    await update.message.reply_document(
                        document=pdf,
                        filename=pdf_filename,
                        caption=f"{caption}\n"
                                f"ðŸ’¸Balansingizdan {config['PDF_COST']} soâ€˜m yechildi.\n"
                                f"ðŸ’µJoriy balans: {user_data['balance']} soâ€˜m\n"
                                f"ðŸ“ŠUmumiy natijalar soni: {user_data['result_count']}",
                        reply_markup=ReplyKeyboardMarkup(
                            [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                             ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                             ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                             ["ðŸ“ž Biz bilan aloqa"]],
                            one_time_keyboard=True, resize_keyboard=True
                        )
                    )
                    logger.info(f"PDF sent to user {user_id}: {pdf_filename}")
            finally:
                # PDF faylini tozalash
                with file_lock:
                    if os.path.exists(output_pdf):
                        try:
                            os.remove(output_pdf)
                            logger.info(f"Removed PDF file: {output_pdf}")
                        except Exception as e:
                            logger.error(f"Error removing PDF file {output_pdf}: {str(e)}")
                    context.user_data.pop('output_pdf', None)

            return START

        except Exception as e:
            logger.error(f"PDF generation failed for user {user_id}: {str(e)}", exc_info=True)
            # Xatolik yuz berganda vaqtinchalik fayllarni tozalash
            with file_lock:
                for file_path in [output_pptx, output_pdf, qr_code_path]:
                    if file_path and os.path.exists(file_path):
                        try:
                            os.remove(file_path)
                            logger.info(f"Removed file: {file_path}")
                        except Exception as e:
                            logger.error(f"Error removing file {file_path}: {str(e)}")
                context.user_data.pop('output_pdf', None)

            await update.message.reply_text(
                f"âŒ PDF faylni tayyorlashda xatolik yuz berdi: {str(e)}",
                reply_markup=ReplyKeyboardMarkup(
                    [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                     ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                     ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                     ["ðŸ“ž Biz bilan aloqa"]],
                    one_time_keyboard=True, resize_keyboard=True
                )
            )
            return START

    await update.message.reply_text(
        "ðŸ“²Iltimos, quyidagi tugmalardan birini tanlang:",
        reply_markup=ReplyKeyboardMarkup([["Ha"], ["Yoâ€˜q"], ["Bosh menyuga qaytish"]], one_time_keyboard=True, resize_keyboard=True)
    )
    return PDF_CONFIRM

async def any_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    if update.message.text == "/start":
        return await start(update, context)
    elif update.message.text == "/admin":
        return await admin_panel(update, context)
    keyboard = [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                ["ðŸ“ž Biz bilan aloqa"]]
    await update.message.reply_text(
        "ðŸ“²Iltimos, quyidagi tugmalardan birini tanlang:",
        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    )
    return START

async def cancel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    keyboard = [["ðŸŸ¢ Sertifikat tayyorlash", "ðŸŸ¢ Taklifnoma yasash"],
                ["ðŸŸ¢ Namuna orqali yasash", "ðŸŸ¢ Diplom yasash"],
                ["ðŸ’° Mening balansim", "ðŸ’¸ Balansni toâ€˜ldirish"],
                ["ðŸ“ž Biz bilan aloqa"]]
    await update.message.reply_text(
        "ðŸ“²Iltimos, quyidagi tugmalardan birini tanlang:",
        reply_markup=ReplyKeyboardMarkup(keyboard, one_time_keyboard=True, resize_keyboard=True)
    )
    return START

async def error_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    logger.error("Xato yuz berdi: %s", str(context.error))
    if update and update.effective_message:
        await update.message.reply_text(
            "âŒ Botda xatolik yuz berdi. Iltimos, keyinroq qayta urinib koâ€˜ring.")

async def start_queue_worker():
    """Queue worker ni ishga tushirish"""
    global queue_worker_running
    if not queue_worker_running:
        asyncio.create_task(queue_worker())

def main():
    logger.info("Bot ishga tushmoqda...")
    config = load_config()
    if not os.path.exists('bot_db.sqlite'):
        logger.warning("bot_db.sqlite fayli yaratilmoqda...")
    try:
        init_db()
        check_fonts()
        logger.info("Ma'lumotlar bazasi muvaffaqiyatli ishga tushdi")
    except Exception as e:
        logger.error(f"Ma'lumotlar bazasini ishga tushirishda xato: {str(e)}")
        return
    app = ApplicationBuilder().token(config['BOT_TOKEN']).build()
    
    # Queue worker ni ishga tushirish
    async def post_init(application):
        await start_queue_worker()
    
    app.post_init = post_init

    app.add_error_handler(error_handler)
    conv_handler = ConversationHandler(
        entry_points=[
            CommandHandler("start", start),
            CommandHandler("admin", admin_panel),
            MessageHandler(filters.TEXT & ~filters.COMMAND, any_message)
        ],
        states={
            START: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_start)],
            DOCUMENT_TYPE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_start)],
            TEMPLATE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_template)],
            SHABLON: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon)],
            SHABLON_NOMI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon_nomi)],
            SHABLON_SHRIFT1: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon_shrift1)],
            SHABLON_MATNI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon_matni)],
            SHABLON_SHRIFT2: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon_shrift2)],
            SHABLON_SHRIFT3: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shablon_shrift3)],
            TAQDIRLANUVCHI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_taqdirlanuvchi)],
            TAKLIFNOMA_SHRIFT1: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_taklifnoma_shrift1)],
            TAKLIFNOMA_SHRIFT2: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_taklifnoma_shrift2)],
            SHRIFT1: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shrift1)],
            TAQDIRLOVCHI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_taqdirlovchi)],
            SERTIFIKAT_MATNI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_sertifikat_matni)],
            DIPLOM_MATNI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_diplom_matni)],
            SHRIFT2: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_shrift2)],
            SANA: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_sana)],
            INFO_TEXT: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_info_text)],
            MANZIL_VA_EGA: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_manzil_va_ega)],
            BALANCE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_balance)],
            PAYMENT_METHOD: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_payment_method)],
            CARD_PAYMENT: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_card_payment)],
            UPLOAD_CHECK: [
                MessageHandler(filters.Document.ALL | filters.PHOTO, handle_upload_check),
                MessageHandler(filters.TEXT & ~filters.COMMAND, handle_upload_check)
            ],
            COMMENT: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_comment)],
            ADMIN_PANEL: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_panel)],
            ADMIN_USER: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_user)],
            ADMIN_ACTION: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_action)],
            ADMIN_FOYDALANUVCHI: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_foydalanuvchi)],
            ADMIN_TOPUP: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_topup)],
            CONFIG_PRICE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_config_price)],
            SET_NEW_PRICE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_set_new_price)],
            ADMIN_MESSAGE_TYPE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_message_type)],
            ADMIN_MESSAGE_RECIPIENT: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_admin_message_recipient)],
            ADMIN_MESSAGE_CONTENT: [
                MessageHandler(filters.TEXT | filters.PHOTO | filters.VIDEO | filters.FORWARDED, handle_admin_message_content)
            ],
            PDF_CONFIRM: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_pdf_confirm)],
            CONTACT: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_contact)],
            CONTACT_MESSAGE: [
                MessageHandler(filters.TEXT | filters.PHOTO, handle_contact_message),
                MessageHandler(filters.TEXT & ~filters.COMMAND, handle_contact_message)
            ],
            QR_CODE: [MessageHandler(filters.TEXT & ~filters.COMMAND, handle_qr_code)]
        },
        fallbacks=[
            CommandHandler("start", start),
            CommandHandler("admin", admin_panel),
            CommandHandler("cancel", cancel)
        ],
    )

    app.add_handler(conv_handler)
    app.add_handler(CallbackQueryHandler(handle_callback_query))

    while True:
        try:
            app.run_polling(
                poll_interval=1.0,
                timeout=30,
                drop_pending_updates=True
            )
        except Exception as e:
            logger.error(f"Pollingda xatolik: {str(e)}")
            time.sleep(2)
            continue


if __name__ == "__main__":
    main()
